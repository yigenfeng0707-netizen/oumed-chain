# -*- coding: utf-8 -*-
"""
Generate 瓯医数链 16-slide presentation (16:9) using python-pptx 1.0.2.

Based on docs/ppt_outline.md (VentureD Hackathon medical track).
Uses colored shapes + text boxes only (no images).

Design system:
  Main  Navy   #1a365d  trust, professional
  Aux   Green  #38a169  health, hope
  Accent Orange #dd6b20 warning, important
  BG    Light  #f7fafc  light gray
  Text  Dark   #2d3748  dark gray
  Font  Microsoft YaHei (Windows)
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

# Design tokens
NAVY      = RGBColor(0x1a, 0x36, 0x5d)
NAVY_DARK = RGBColor(0x10, 0x24, 0x40)
NAVY_LT   = RGBColor(0x2c, 0x4f, 0x7c)
GREEN     = RGBColor(0x38, 0xa1, 0x69)
GREEN_LT  = RGBColor(0xc6, 0xe6, 0xd5)
GREEN_BG  = RGBColor(0xed, 0xf7, 0xf0)
ORANGE    = RGBColor(0xdd, 0x6b, 0x20)
ORANGE_LT = RGBColor(0xfe, 0xee, 0xe2)
ORANGE_BG = RGBColor(0xff, 0xf5, 0xeb)
LIGHT_BG  = RGBColor(0xf7, 0xfa, 0xfc)
WHITE     = RGBColor(0xff, 0xff, 0xff)
DARK      = RGBColor(0x2d, 0x37, 0x48)
GRAY      = RGBColor(0x71, 0x80, 0x96)
GRAY_LT   = RGBColor(0xe2, 0xe8, 0xf0)
RED       = RGBColor(0xc5, 0x30, 0x30)
YELLOW    = RGBColor(0xd6, 0x9e, 0x2e)
SUB_CL    = RGBColor(0xcb, 0xd5, 0xe0)
PAGE_CL   = RGBColor(0xa0, 0xae, 0xc0)

FONT = "Microsoft YaHei"

SW = 13.333
SH = 7.5


def _font(run, size, color, bold, name=FONT):
    run.font.name = name
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    rPr = run._r.get_or_add_rPr()
    for tag in ('a:ea', 'a:cs'):
        el = rPr.find(qn(tag))
        if el is None:
            el = rPr.makeelement(qn(tag), {})
            rPr.append(el)
        el.set('typeface', name)


def set_text(tf, lines, size=14, color=DARK, bold=False, align=PP_ALIGN.LEFT,
             anchor=MSO_ANCHOR.TOP):
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = Pt(5); tf.margin_right = Pt(5)
    tf.margin_top = Pt(2);  tf.margin_bottom = Pt(2)
    if isinstance(lines, str):
        lines = [lines]
    for i, item in enumerate(lines):
        if isinstance(item, str):
            t, sz, cl, bd, al = item, size, color, bold, align
        else:
            t  = item[0]
            sz = item[1] if len(item) > 1 else size
            cl = item[2] if len(item) > 2 else color
            bd = item[3] if len(item) > 3 else bold
            al = item[4] if len(item) > 4 else align
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = al
        r = p.add_run()
        r.text = t
        _font(r, sz, cl, bd)


def rect(s, l, t, w, h, fill, line=None, lw=1.0):
    sp = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(l), Inches(t), Inches(w), Inches(h))
    sp.fill.solid(); sp.fill.fore_color.rgb = fill
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line; sp.line.width = Pt(lw)
    sp.shadow.inherit = False
    return sp


def rrect(s, l, t, w, h, fill, line=None, lw=1.0, radius=0.06):
    sp = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(l), Inches(t), Inches(w), Inches(h))
    sp.fill.solid(); sp.fill.fore_color.rgb = fill
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line; sp.line.width = Pt(lw)
    try:
        sp.adjustments[0] = radius
    except Exception:
        pass
    sp.shadow.inherit = False
    return sp


def oval(s, l, t, d, fill, line=None, lw=1.0):
    sp = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(l), Inches(t), Inches(d), Inches(d))
    sp.fill.solid(); sp.fill.fore_color.rgb = fill
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line; sp.line.width = Pt(lw)
    sp.shadow.inherit = False
    return sp


def tbox(s, l, t, w, h, lines, **kw):
    bx = s.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    set_text(bx.text_frame, lines, **kw)
    return bx


def icon(s, l, t, d, fill, label, lc=WHITE, size=14):
    c = oval(s, l, t, d, fill)
    set_text(c.text_frame, [label], size=size, color=lc, bold=True,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    return c


def set_bg(s, color):
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = color


def notes(s, text):
    s.notes_slide.notes_text_frame.text = text


def title_bar(s, title, subtitle=None, page=None, accent=GREEN):
    rect(s, 0, 0, SW, 1.0, NAVY)
    rect(s, 0, 1.0, SW, 0.06, accent)
    tbox(s, 0.5, 0.08, 11.4, 0.55, [title], size=24, color=WHITE, bold=True,
         anchor=MSO_ANCHOR.MIDDLE)
    if subtitle:
        tbox(s, 0.5, 0.60, 11.4, 0.36, [subtitle], size=12, color=SUB_CL,
             anchor=MSO_ANCHOR.MIDDLE)
    if page:
        tbox(s, 11.9, 0.30, 1.1, 0.4, [str(page) + " / 16"], size=11, color=PAGE_CL,
             align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE)


def card(s, l, t, w, h, fill=WHITE, line=None, lw=1.0, radius=0.06):
    return rrect(s, l, t, w, h, fill, line, lw, radius)


def card_title(s, l, t, w, title, fill=NAVY, tc=WHITE, size=14):
    bar = rect(s, l, t, w, 0.42, fill)
    set_text(bar.text_frame, [title], size=size, color=tc, bold=True,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    return bar


# ---- Slide builders ----
def slide01(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s, NAVY)
    rect(s, 0, 0, SW, 0.22, GREEN)
    rect(s, 0, SH - 0.22, SW, 0.22, NAVY_DARK)
    tbox(s, 1.0, 1.55, 11.3, 1.0, ["瓯医数链"], size=54, color=WHITE, bold=True,
         align=PP_ALIGN.CENTER)
    tbox(s, 1.0, 2.60, 11.3, 0.5, ["multi-modal medical signal agent"], size=14,
         color=SUB_CL, align=PP_ALIGN.CENTER)
    tbox(s, 1.0, 3.05, 11.3, 0.45, [("\u591a\u6a21\u6001\u533b\u7597\u4fe1\u53f7\u667a\u80fd\u4f53", 22, SUB_CL, True, PP_ALIGN.CENTER)])
    tbox(s, 1.0, 3.55, 11.3, 0.4,
         [("\u5173\u952e\u533b\u7597\u4fe1\u53f7\u8bc6\u522b \u00d7 \u60a3\u8005\u4fe1\u606f\u8fde\u63a5", 16, GREEN_LT, True, PP_ALIGN.CENTER)])
    sigs = [("\u8111\u6ce2", GREEN), ("\u5f71\u50cf", ORANGE), ("\u5fc3\u8df3", GREEN)]
    cw, gap = 1.5, 0.45
    total = cw * 3 + gap * 2
    x0 = (SW - total) / 2
    for i, (lab, col) in enumerate(sigs):
        icon(s, x0 + i * (cw + gap), 4.1, cw, col, lab, size=16)
    tbox(s, 0.5, 5.65, 12.3, 0.3,
         [("\u4e09\u5927\u4fe1\u53f7 \u00b7 \u8111\u7535 / \u5f71\u50cf / \u884c\u4e3a", 12, SUB_CL, False, PP_ALIGN.CENTER)])
    agents = [("\u8111\u7535", GREEN), ("\u5f71\u50cf", GREEN), ("\u5065\u5eb7", GREEN),
              ("\u6743\u76ca", NAVY_LT), ("\u62a5\u9500", NAVY_LT), ("\u653f\u7b56", NAVY_LT),
              ("\u5b89\u5168", ORANGE)]
    n = len(agents)
    aw, g2 = 0.95, 0.22
    tot = aw * n + g2 * (n - 1)
    st = (SW - tot) / 2
    for i, (lab, col) in enumerate(agents):
        icon(s, st + i * (aw + g2), 6.1, aw, col, lab, size=10)
    tbox(s, 0.5, 7.1, 12.3, 0.3,
         [("VentureD Hackathon \u00b7 \u533b\u7597\u8d5b\u9053      |      2026 \u5e74 8 \u6708", 12, SUB_CL, False, PP_ALIGN.CENTER)])
    notes(s, "\u5404\u4f4d\u8bc4\u59d4\u3001\u5404\u4f4d\u8001\u5e08\uff0c\u5927\u5bb6\u597d\u3002\u4eca\u5929\u6211\u5e26\u6765\u7684\u662f 瓯医数链\u2014\u2014\u4e00\u4e2a\u9762\u5411\u771f\u5b9e\u533b\u7597\u573a\u666f\u3001\u4e13\u95e8\u8bc6\u522b\u201c\u5173\u952e\u533b\u7597\u4fe1\u53f7\u201d\u7684 AI \u667a\u80fd\u4f53\u3002")


def slide02(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s, LIGHT_BG)
    title_bar(s, "\u88ab\u9519\u8fc7\u7684\u5173\u952e\u533b\u7597\u4fe1\u53f7",
              "\u4e09\u7c7b\u771f\u5b9e\u573a\u666f\uff0c\u4e09\u7c7b\u88ab\u5ffd\u89c6\u7684\u533b\u7597\u4fe1\u53f7", 2, accent=ORANGE)
    cases = [
        ("\u5f20\u5148\u751f \u00b7 55 \u5c81 \u00b7 \u957f\u671f\u52a0\u73ed",
         "\u538b\u529b\u4e0e\u7761\u7720\u95ee\u9898\u6ca1\u6709\u5ba2\u89c2\u6307\u6807\uff0c\u76f4\u5230\u4f53\u68c0\u62a5\u544a\u4eae\u7ea2\u706f",
         "\u8111\u7535\u4fe1\u53f7\uff1a\u538b\u529b/\u7761\u7720\u662f\u770b\u4e0d\u89c1\u7684\u5065\u5eb7\u4fe1\u53f7", GREEN, "\u8111"),
        ("\u57fa\u5c42\u533b\u751f \u00b7 \u65e5\u5747\u9605\u7247\u4e0a\u767e\u5f20",
         "\u5de5\u4f5c 8 \u5c0f\u65f6\u540e\u6f0f\u8bca\u98ce\u9669\u663e\u8457\u4e0a\u5347\uff0c\u75c5\u7076\u85cf\u5728\u5f71\u50cf\u91cc",
         "\u5f71\u50cf\u4fe1\u53f7\uff1a\u75c5\u7076\u662f\u85cf\u4e0d\u4f4f\u7684\u89c6\u89c9\u4fe1\u53f7", ORANGE, "\u5f71"),
        ("\u666e\u901a\u60a3\u8005 \u00b7 \u9762\u5bf9\u6ee1\u7eb8\u653f\u7b56",
         "\u770b\u4e0d\u61c2\u653f\u7b56\u4e0e\u533b\u7597\u4fe1\u606f\uff0c\u8be5\u4eab\u53d7\u7684\u5f85\u9047\u4e00\u518d\u9519\u8fc7",
         "\u4fe1\u606f\u4fe1\u53f7\uff1a\u653f\u7b56\u7ea2\u5229\u662f\u62ff\u4e0d\u5230\u7684\u4ef7\u503c\u4fe1\u53f7", NAVY_LT, "\u4fe1"),
    ]
    cw, gap = 3.9, 0.3
    x0 = 0.5
    for i, (head, body, sig, col, lab) in enumerate(cases):
        x = x0 + i * (cw + gap)
        card(s, x, 1.35, cw, 3.45, WHITE, line=GRAY_LT, lw=1.0)
        icon(s, x + cw / 2 - 0.4, 1.5, 0.8, col, lab, size=18)
        tbox(s, x + 0.2, 2.42, cw - 0.4, 0.4, [head], size=13, color=NAVY, bold=True,
             align=PP_ALIGN.CENTER)
        tbox(s, x + 0.2, 2.85, cw - 0.4, 1.0, [body], size=11, color=DARK,
             align=PP_ALIGN.CENTER)
        sig_bar = rrect(s, x + 0.15, 4.05, cw - 0.3, 0.62, col, radius=0.18)
        set_text(sig_bar.text_frame, [sig], size=11, color=WHITE, bold=True,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    call = rrect(s, 0.5, 5.05, 12.33, 0.95, NAVY, radius=0.10)
    set_text(call.text_frame,
             [("\u5173\u952e\u533b\u7597\u4fe1\u53f7\uff0c\u6b63\u5728\u88ab\u9519\u8fc7", 28, ORANGE, True, PP_ALIGN.CENTER)],
             anchor=MSO_ANCHOR.MIDDLE)
    notes(s, "\u5f20\u5148\u751f 55 \u5c81\uff0c\u957f\u671f\u52a0\u73ed\uff0c\u538b\u529b\u5927\u3001\u7761\u4e0d\u597d\uff0c\u53ef\u8fd9\u4e9b\u90fd\u6ca1\u6709\u5ba2\u89c2\u6307\u6807\uff0c\u76f4\u5230\u4f53\u68c0\u62a5\u544a\u4eae\u7ea2\u706f\u3002\u53e6\u4e00\u8fb9\uff0c\u57fa\u5c42\u533b\u751f\u6bcf\u5929\u8981\u9605\u4e0a\u767e\u5f20\u7247\u5b50\uff0c8 \u5c0f\u65f6\u540e\u7684\u6f0f\u8bca\u98ce\u9669\u663e\u8457\u4e0a\u5347\u3002\u800c\u6211\u4eec\u7684\u60a3\u8005\uff0c\u9762\u5bf9\u6ee1\u7eb8\u653f\u7b56\u6587\u4ef6\uff0c\u4e0d\u77e5\u9053\u8be5\u4eab\u53d7\u4ec0\u4e48\u5f85\u9047\u3002\u4e09\u7c7b\u5173\u952e\u533b\u7597\u4fe1\u53f7\u2014\u2014\u8111\u7535\u7684\u3001\u5f71\u50cf\u7684\u3001\u4fe1\u606f\u7684\u2014\u2014\u6b63\u5728\u88ab\u9519\u8fc7\u3002")


def slide03(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s, LIGHT_BG)
    title_bar(s, "\u533b\u7597\u4fe1\u53f7\u8bc6\u522b\u7684\u5de8\u5927\u7f3a\u53e3",
             "\u4e09\u7ec4\u6570\u636e\uff0c\u540c\u4e00\u4e2a\u95ee\u9898", 3, accent=ORANGE)
    data = [
        ("\u8111\u7535", "\u6301\u7eed\u8d70\u9ad8", "\u4e2d\u56fd\u6210\u5e74\u4eba\u6162\u6027\u75c5\u60a3\u75c5\u7387\u6301\u7eed\u8d70\u9ad8\uff0c\u5fc3\u7406\u5065\u5eb7\u4e0e\u7761\u7720\u969c\u788d\u7f3a\u4e4f\u53ef\u53ca\u3001\u5ba2\u89c2\u7684\u65e9\u671f\u7b5b\u67e5\u624b\u6bb5", GREEN, "\u8111"),
        ("\u5f71\u50cf", "\u4f9b\u9700\u5931\u8861", "\u57fa\u5c42\u5f71\u50cf\u533b\u5e08\u7f3a\u53e3\u5927\uff0c\u65e5\u5747\u9605\u7247\u91cf\u8d85\u8d1f\u8377\uff0c\u65e9\u671f\u75c5\u7076\u6f0f\u68c0\u7387\u4e0d\u5bb9\u5ffd\u89c6", ORANGE, "\u5f71"),
        ("\u4fe1\u606f", "\u7ea2\u5229\u6d41\u5931", "\u5927\u91cf\u53c2\u4fdd\u4eba\u4e0d\u77e5\u9053\u81ea\u5df1\u80fd\u4eab\u53d7\u54ea\u4e9b\u533b\u7597\u5f85\u9047\uff0c\u6bcf\u5e74\u9519\u8fc7\u7684\u653f\u7b56\u7ea2\u5229\u4ee5\u5343\u5143\u8ba1", NAVY_LT, "\u4fe1"),
    ]
    cw, gap = 3.9, 0.3
    x0 = 0.5
    for i, (tag, big, desc, col, lab) in enumerate(data):
        x = x0 + i * (cw + gap)
        card(s, x, 1.4, cw, 4.3, WHITE, line=GRAY_LT, lw=1.0)
        rect(s, x, 1.4, cw, 0.5, col)
        tbox(s, x, 1.42, cw, 0.46, [(tag + "\u4fe1\u53f7", 14, WHITE, True, PP_ALIGN.CENTER)],
             anchor=MSO_ANCHOR.MIDDLE)
        icon(s, x + cw / 2 - 0.45, 2.05, 0.9, col, lab, size=20)
        tbox(s, x + 0.2, 3.05, cw - 0.4, 0.7, [(big, 26, ORANGE, True, PP_ALIGN.CENTER)])
        tbox(s, x + 0.25, 3.85, cw - 0.5, 1.7, [desc], size=11.5, color=DARK,
             align=PP_ALIGN.CENTER)
    foot = rrect(s, 0.5, 6.05, 12.33, 0.7, NAVY, radius=0.10)
    set_text(foot.text_frame,
             [("\u4e13\u4e1a\u4fe1\u53f7\u4e0e\u666e\u901a\u4eba\u4e4b\u95f4\u5b58\u5728\u5de8\u5927\u9e3f\u6c9f \u2014\u2014 \u8bc6\u522b\u4e0e\u7ffb\u8bd1\uff0c\u6b63\u662f\u6211\u4eec\u7684\u5207\u5165\u70b9", 14, WHITE, True, PP_ALIGN.CENTER)],
             anchor=MSO_ANCHOR.MIDDLE)
    notes(s, "\u4e09\u7ec4\u6570\u636e\u8bf4\u660e\u540c\u4e00\u4e2a\u95ee\u9898\uff1a\u65e0\u8bba\u8111\u7535\u3001\u5f71\u50cf\u8fd8\u662f\u653f\u7b56\u4fe1\u606f\uff0c\u4e13\u4e1a\u4fe1\u53f7\u4e0e\u666e\u901a\u4eba\u4e4b\u95f4\u5b58\u5728\u5de8\u5927\u9e3f\u6c9f\u3002\u533b\u7597\u4fe1\u53f7\u7684\u8bc6\u522b\u4e0e\u7ffb\u8bd1\uff0c\u6b63\u662f\u6211\u4eec\u5207\u5165\u7684\u771f\u5b9e\u573a\u666f\u3002")


def slide04(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s, LIGHT_BG)
    title_bar(s, "瓯医数链 \u2014\u2014 \u8ba9\u5173\u952e\u533b\u7597\u4fe1\u53f7\uff0c\u4e0d\u518d\u88ab\u9519\u8fc7",
             "\u4e00\u4e2a\u80fd\u201c\u770b\u4fe1\u53f7\u3001\u6293\u5f02\u5e38\u3001\u8fde\u8d44\u6e90\u201d\u7684\u591a\u6a21\u6001\u533b\u7597\u4fe1\u53f7\u667a\u80fd\u4f53", 4, accent=GREEN)
    pos = rrect(s, 2.0, 1.35, 9.33, 0.95, NAVY, radius=0.10)
    set_text(pos.text_frame,
             [("\u4e00\u53e5\u8bdd\u5b9a\u4f4d\uff1a\u8bc6\u522b\u4fe1\u53f7 \u00b7 \u6293\u53d6\u5f02\u5e38 \u00b7 \u8fde\u63a5\u8d44\u6e90", 18, WHITE, True, PP_ALIGN.CENTER)],
             anchor=MSO_ANCHOR.MIDDLE)
    caps = [
        ("\u8bc6\u522b\u8111\u7535\u4fe1\u53f7", "EEG \u9891\u57df\u5206\u6790", "\u63a5\u5165 LSL \u517c\u5bb9\u8bbe\u5907\nFFT + Welch PSD\n\u8f93\u51fa\u4e94\u7ef4\u5065\u5eb7\u6307\u6807", GREEN, "\u8111"),
        ("\u8bc6\u522b\u5f71\u50cf\u4fe1\u53f7", "AI \u9884\u6807\u6ce8 + \u533b\u5e08\u590d\u6838", "\u75c5\u7076\u68c0\u6d4b bbox \u6846\u9009\n\u533b\u5e08\u9010\u6846\u786e\u8ba4/\u9a73\u56de\n\u7ed3\u6784\u5316\u62a5\u544a\u95ed\u73af", ORANGE, "\u5f71"),
        ("\u8fde\u63a5\u60a3\u8005\u4fe1\u606f", "\u653f\u7b56 / \u62a5\u9500 / \u6743\u76ca", "\u653f\u7b56\u5339\u914d\u6309\u7701\u94b1\u6392\u5e8f\n\u62a5\u9500\u9884\u5ba1 7 \u6b65\u63a8\u5bfc\n\u6743\u76ca\u95ee\u7b54\u4e00\u53e5\u8bdd\u67e5\u6e05", NAVY_LT, "\u4fe1"),
    ]
    cw, gap = 3.9, 0.3
    x0 = 0.5
    for i, (head, sub, body, col, lab) in enumerate(caps):
        x = x0 + i * (cw + gap)
        card(s, x, 2.55, cw, 3.5, WHITE, line=col, lw=1.5)
        icon(s, x + cw / 2 - 0.45, 2.7, 0.9, col, lab, size=20)
        tbox(s, x + 0.2, 3.72, cw - 0.4, 0.4, [head], size=14, color=NAVY, bold=True,
             align=PP_ALIGN.CENTER)
        tbox(s, x + 0.2, 4.12, cw - 0.4, 0.32, [sub], size=11, color=col, bold=True,
             align=PP_ALIGN.CENTER)
        tbox(s, x + 0.25, 4.55, cw - 0.5, 1.4, [body], size=11, color=DARK,
             align=PP_ALIGN.CENTER)
    slogan = rrect(s, 0.5, 6.25, 12.33, 0.7, GREEN, radius=0.10)
    set_text(slogan.text_frame,
             [("\u8bc6\u522b\u4fe1\u53f7 \u00b7 \u5b88\u62a4\u5065\u5eb7 \u00b7 \u8fde\u63a5\u8d44\u6e90", 18, WHITE, True, PP_ALIGN.CENTER)],
             anchor=MSO_ANCHOR.MIDDLE)
    notes(s, "\u6240\u4ee5\u6211\u4eec\u505a\u4e86 瓯医数链\u3002\u5b83\u4e0d\u53ea\u56de\u7b54\u95ee\u9898\uff0c\u5b83\u80fd\u8bc6\u522b\u8111\u7535\u4fe1\u53f7\u3001\u6293\u53d6\u5f71\u50cf\u75c5\u7076\u3001\u8fde\u63a5\u653f\u7b56\u4e0e\u62a5\u9500\u8d44\u6e90\uff0c\u628a\u201c\u770b\u4e0d\u89c1\u7684\u5f02\u5e38\u201d\u53d8\u6210\u201c\u770b\u5f97\u89c1\u7684\u884c\u52a8\u201d\u3002")


def slide05(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s, LIGHT_BG)
    title_bar(s, "\u591a\u6a21\u6001\u4fe1\u53f7\u8bc6\u522b\u67b6\u6784",
             "\u7f16\u6392\u667a\u80fd\u4f53 \u00d7 7 \u4e2a\u4e13\u4e1a\u667a\u80fd\u4f53 \u00d7 \u5b89\u5168\u5408\u89c4\u5c42", 5, accent=GREEN)
    orch = rrect(s, 0.5, 1.3, 12.33, 0.75, NAVY, radius=0.10)
    set_text(orch.text_frame,
             [("\u7f16\u6392\u667a\u80fd\u4f53 Orchestrator \u2014\u2014 \u610f\u56fe\u8bc6\u522b \u00b7 \u4fe1\u53f7\u8def\u7531 \u00b7 \u7ed3\u679c\u805a\u5408", 15, WHITE, True, PP_ALIGN.CENTER)],
             anchor=MSO_ANCHOR.MIDDLE)
    tbox(s, 0.5, 2.18, 6.0, 0.3, [("\u4fe1\u53f7\u8bc6\u522b\u5c42 \u2b50 \u6838\u5fc3", 12, GREEN, True)])
    sig_agents = [
        ("\u8111\u7535\u536b\u58eb", "EEG \u9891\u57df\u5206\u6790\n\u4e94\u7ef4\u5065\u5eb7\u6307\u6807", GREEN, True),
        ("\u5f71\u50cf\u536b\u58eb", "\u75c5\u7076\u68c0\u6d4b \u2192 \u9884\u6807\u6ce8\n\u533b\u5e08\u590d\u6838 \u2192 \u62a5\u544a", ORANGE, True),
        ("\u5065\u5eb7\u536b\u58eb", "\u7528\u836f/\u5c31\u533b\u6a21\u5f0f\n\u6162\u75c5\u98ce\u9669\u9884\u8b66", GREEN, False),
    ]
    cw, gap = 3.9, 0.3
    x0 = 0.5
    for i, (name, body, col, star) in enumerate(sig_agents):
        x = x0 + i * (cw + gap)
        line_c = col if star else GRAY_LT
        card(s, x, 2.5, cw, 1.45, WHITE, line=line_c, lw=2.0 if star else 1.0)
        tbox(s, x + 0.2, 2.58, cw - 0.4, 0.35,
             [(name + (" \u2b50" if star else ""), 13, col if star else NAVY, True, PP_ALIGN.CENTER)])
        tbox(s, x + 0.2, 2.95, cw - 0.4, 0.95, [body], size=10.5, color=DARK,
             align=PP_ALIGN.CENTER)
    tbox(s, 0.5, 4.08, 6.0, 0.3, [("\u4fe1\u606f\u8fde\u63a5\u5c42", 12, NAVY_LT, True)])
    info_agents = [
        ("\u6743\u76ca\u7ba1\u5bb6", "\u53c2\u4fdd\u7c7b\u578b/\u8d26\u6237\u4f59\u989d/\u62a5\u9500\u6bd4\u4f8b", NAVY_LT),
        ("\u62a5\u9500\u52a9\u624b", "\u7968\u636e OCR + 7 \u6b65\u62a5\u9500\u8ba1\u7b97", NAVY_LT),
        ("\u653f\u7b56\u53c2\u8c0b", "\u6309\u7701\u94b1\u91d1\u989d\u6392\u5e8f + \u539f\u6587\u5f15\u7528", NAVY_LT),
    ]
    for i, (name, body, col) in enumerate(info_agents):
        x = x0 + i * (cw + gap)
        card(s, x, 4.4, cw, 1.15, WHITE, line=GRAY_LT, lw=1.0)
        tbox(s, x + 0.2, 4.46, cw - 0.4, 0.35, [(name, 13, NAVY, True, PP_ALIGN.CENTER)])
        tbox(s, x + 0.2, 4.82, cw - 0.4, 0.65, [body], size=10.5, color=DARK,
             align=PP_ALIGN.CENTER)
    sec = rrect(s, 0.5, 5.7, 12.33, 0.65, ORANGE, radius=0.10)
    set_text(sec.text_frame,
             [("\u5b89\u5168\u5408\u89c4\u5c42 \u00b7 \u5b89\u5168\u5b88\u95e8 \u2014\u2014 \u6388\u6743\u7ba1\u7406 + \u53ef\u4fe1\u6570\u636e\u7a7a\u95f4 + \u533a\u5757\u94fe\u5b58\u8bc1", 14, WHITE, True, PP_ALIGN.CENTER)],
             anchor=MSO_ANCHOR.MIDDLE)
    tbox(s, 0.5, 6.48, 12.33, 0.4,
         [("用户输入 → 编排器 → 多智能体并行 → 结果聚合 → 用户 · 会话持久化 · 动态用户全站联动", 11, GRAY, False, PP_ALIGN.CENTER)])
    notes(s, "\u67b6\u6784\u5206\u4e09\u5c42\uff1a\u7f16\u6392\u5668\u8d1f\u8d23\u8c03\u5ea6\uff0c\u4fe1\u53f7\u8bc6\u522b\u5c42\u8986\u76d6\u8111\u7535\u3001\u5f71\u50cf\u3001\u884c\u4e3a\u4e09\u7c7b\u5173\u952e\u4fe1\u53f7\uff0c\u4fe1\u606f\u8fde\u63a5\u5c42\u628a\u7ed3\u679c\u7ffb\u8bd1\u6210\u60a3\u8005\u542c\u5f97\u61c2\u3001\u7528\u5f97\u4e0a\u7684\u5efa\u8bae\uff0c\u5e95\u5c42\u7531\u5b89\u5168\u5b88\u95e8\u4fdd\u969c\u9690\u79c1\u5408\u89c4\u3002\u5176\u4e2d\u8111\u7535\u536b\u58eb\u548c\u5f71\u50cf\u536b\u58eb\u662f\u6211\u4eec\u7684\u6838\u5fc3\u521b\u65b0\u3002")


def slide06(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s, LIGHT_BG)
    title_bar(s, "Demo 1 \u00b7 \u8111\u7535\u4fe1\u53f7\u8bc6\u522b \u2014\u2014 \u770b\u4e0d\u89c1\u7684\u538b\u529b\u6709\u4e86\u5ba2\u89c2\u6307\u6807",
             "EEG \u9891\u57df\u5206\u6790 \u00b7 \u4e94\u7ef4\u5065\u5eb7\u6307\u6807 \u00b7 \u4e3b\u52a8\u9884\u8b66", 6, accent=GREEN)
    card(s, 0.5, 1.35, 4.3, 4.4, WHITE, line=GRAY_LT, lw=1.0)
    card_title(s, 0.5, 1.35, 4.3, "EEG \u8bbe\u5907\u63a5\u5165", fill=GREEN)
    specs = ["Muse / Emotiv / OpenBCI", "4 通道 · 256 Hz 采样",
             "LSL 兼容 · 仿真信号演示", "FFT + Welch PSD 频域分析",
             "五频段 δ / θ / α / β / γ", "双源真实数据集验证"]
    tbox(s, 0.7, 1.95, 3.9, 2.4, [(ln, 11.5, DARK) for ln in specs])
    wf = rect(s, 0.7, 4.45, 3.9, 1.15, NAVY_DARK)
    tbox(s, 0.7, 4.45, 3.9, 1.15, [("\u5b9e\u65f6\u6ce2\u5f62 \u00b7 4 \u901a\u9053\u6eda\u52a8", 10, GREEN_LT, False, PP_ALIGN.CENTER)],
         anchor=MSO_ANCHOR.MIDDLE)
    card(s, 5.0, 1.35, 4.3, 4.4, WHITE, line=GRAY_LT, lw=1.0)
    card_title(s, 5.0, 1.35, 4.3, "\u4e94\u7ef4\u5065\u5eb7\u6307\u6807\u8f93\u51fa", fill=NAVY)
    metrics = [("\u538b\u529b\u6307\u6570", "68", ORANGE), ("\u6ce8\u610f\u529b", "42", NAVY_LT),
               ("\u7761\u7720\u8d28\u91cf", "55", GREEN), ("\u8ba4\u77e5\u8d1f\u8377", "71", ORANGE),
               ("\u60c5\u7eea\u72b6\u6001", "49", NAVY_LT)]
    my = 1.95
    for name, val, col in metrics:
        rect(s, 5.2, my, 3.9, 0.62, LIGHT_BG)
        tbox(s, 5.35, my, 2.6, 0.62, [(name, 12, DARK, False)], anchor=MSO_ANCHOR.MIDDLE)
        tbox(s, 7.9, my, 1.0, 0.62, [(val, 16, col, True, PP_ALIGN.RIGHT)], anchor=MSO_ANCHOR.MIDDLE)
        my += 0.72
    card(s, 9.5, 1.35, 3.33, 4.4, WHITE, line=ORANGE, lw=1.5)
    card_title(s, 9.5, 1.35, 3.33, "\u4e3b\u52a8\u9884\u8b66", fill=ORANGE)
    icon(s, 10.83, 2.0, 1.0, ORANGE, "!", size=22)
    tbox(s, 9.7, 3.1, 2.93, 0.9, [("\u68c0\u6d4b\u5230\u9ad8\u538b\u529b + \u7761\u7720\u4e0d\u8db3", 13, ORANGE, True, PP_ALIGN.CENTER)])
    tbox(s, 9.7, 4.05, 2.93, 0.5, [("\u2192 \u5339\u914d\u533b\u7597\u5f85\u9047", 11, DARK, False, PP_ALIGN.CENTER)])
    pol = rrect(s, 9.7, 4.6, 2.93, 0.95, GREEN_BG, line=GREEN, lw=1.0, radius=0.10)
    set_text(pol.text_frame, [("\u7761\u7720\u969c\u788d\u95e8\u8bca\u5f85\u9047\n\u653f\u7b56\u539f\u6587\u5f15\u7528", 10.5, GREEN, True, PP_ALIGN.CENTER)],
             anchor=MSO_ANCHOR.MIDDLE)
    badge = rrect(s, 0.5, 6.0, 12.33, 0.75, GREEN, radius=0.10)
    set_text(badge.text_frame,
             [("✓ 双源真实数据集 · PhysioNet 运动想象 + EEGEmotions-27 情绪 · 289 项测试全绿", 15, WHITE, True, PP_ALIGN.CENTER)],
             anchor=MSO_ANCHOR.MIDDLE)
    notes(s, "第一个 Demo。接入 EEG 设备，瓯医数链 对脑电做频域分析，输出压力、睡眠等五个维度的客观指标。张先生的“睡不好”第一次有了数字。当指标异常，它主动预警，并自动匹配到相关医疗待遇。引擎已在双源真实公开数据集上验证：PhysioNet 运动想象数据集 + EEGEmotions-27 情绪数据集，加上 289 项单元测试在支撑。")


def slide07(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s, ORANGE_BG)
    title_bar(s, "Demo 2 \u00b7 \u533b\u5b66\u5f71\u50cf\u8bc6\u522b \u2014\u2014 AI \u9884\u6807\u6ce8\uff0c\u533b\u5e08\u505a\u88c1\u5224",
             "\u2b50 KILLER FEATURE \u00b7 \u4eba\u673a\u534f\u540c\u95ed\u73af", 7, accent=ORANGE)
    rect(s, 0.35, 1.25, 12.63, 4.7, ORANGE_BG)
    rect(s, 0.35, 1.25, 0.12, 4.7, ORANGE)
    rect(s, 0.35, 1.25, 12.63, 0.12, ORANGE)
    scan = rect(s, 0.7, 1.55, 5.6, 4.05, GRAY_LT)
    tbox(s, 0.7, 1.55, 5.6, 0.4, [("\u80f8\u7247 / \u80ba CT / \u8111 MRI", 11, GRAY, False, PP_ALIGN.CENTER)],
         anchor=MSO_ANCHOR.MIDDLE)
    rrect(s, 1.5, 2.3, 1.7, 1.2, WHITE, line=RED, lw=2.5, radius=0.0)
    tbox(s, 1.5, 2.05, 1.7, 0.25, [("\u9ad8\u5371 0.92", 9, RED, True, PP_ALIGN.CENTER)])
    rrect(s, 3.6, 3.0, 1.4, 1.0, WHITE, line=YELLOW, lw=2.5, radius=0.0)
    tbox(s, 3.6, 2.78, 1.4, 0.22, [("\u4e2d\u5371 0.74", 9, YELLOW, True, PP_ALIGN.CENTER)])
    rrect(s, 2.6, 4.35, 1.1, 0.8, WHITE, line=GREEN, lw=2.5, radius=0.0)
    tbox(s, 2.6, 4.15, 1.1, 0.2, [("\u4f4e\u5371 0.61", 9, GREEN, True, PP_ALIGN.CENTER)])
    tbox(s, 0.8, 5.05, 5.4, 0.5, [("\u200bAI \u68c0\u51fa 3 \u5904\u75c5\u7076 \u00b7 bbox \u6309\u4e25\u91cd\u5ea6\u7740\u8272", 10, DARK, True, PP_ALIGN.CENTER)])
    card(s, 6.55, 1.55, 6.1, 4.05, WHITE, line=GRAY_LT, lw=1.0)
    card_title(s, 6.55, 1.55, 6.1, "\u533b\u5e08\u9010\u6846\u590d\u6838\u9762\u677f", fill=NAVY)
    steps = [("1", "\u9009\u62e9\u68c0\u67e5\u7c7b\u578b \u2192 \u4e00\u952e AI \u5206\u6790"),
             ("2", "AI \u68c0\u51fa\u75c5\u7076\u5e76 bbox \u6846\u9009\uff08\u9ad8\u5371\u7ea2/\u4e2d\u5371\u9ec4/\u4f4e\u5371\u7eff\uff09"),
             ("3", "\u533b\u5e08\u9010\u6846\u64cd\u4f5c\uff1a\u786e\u8ba4 / \u9a73\u56de / \u70b9\u51fb\u5f71\u50cf\u65b0\u589e\u6807\u6ce8"),
             ("4", "\u590d\u6838\u5b8c\u6210 \u2192 \u751f\u6210\u7ed3\u6784\u5316\u62a5\u544a\uff08\u75c5\u7076\u6e05\u5355 + \u98ce\u9669\u7b49\u7ea7 + \u5efa\u8bae\uff09"),
             ("5", "\u5f71\u50cf\u5f02\u5e38\u81ea\u52a8\u8054\u52a8\u76f8\u5173\u653f\u7b56\u4e0e\u5c31\u8bca\u5efa\u8bae")]
    sy = 2.12
    for n, txt in steps:
        icon(s, 6.75, sy, 0.42, NAVY, n, size=12)
        tbox(s, 7.3, sy, 5.2, 0.6, [(txt, 10.5, DARK)], anchor=MSO_ANCHOR.MIDDLE)
        sy += 0.62
    tbox(s, 6.75, 5.15, 5.7, 0.4, [("\u2713 \u5f71\u50cf\u5f15\u64ce\u6d4b\u8bd5\u901a\u8fc7 \u00b7 bbox \u5750\u6807\u5f52\u4e00\u5316\u7cbe\u5ea6\u9a8c\u8bc1", 9.5, GREEN, True, PP_ALIGN.CENTER)])
    call = rrect(s, 0.5, 6.1, 12.33, 0.85, ORANGE, radius=0.10)
    set_text(call.text_frame,
             [("AI \u9884\u6807\u6ce8\u662f\u5efa\u8bae\uff0c\u533b\u5e08\u590d\u6838\u662f\u88c1\u51b3 \u2014\u2014 \u533b\u7597 AI \u5e94\u6709\u7684\u5b89\u5168\u8fb9\u754c", 17, WHITE, True, PP_ALIGN.CENTER)],
             anchor=MSO_ANCHOR.MIDDLE)
    notes(s, "\u7b2c\u4e8c\u4e2a Demo \u662f\u5f71\u50cf\u8bc6\u522b\u3002AI \u5728 CT \u4e0a\u81ea\u52a8\u6846\u51fa\u75c5\u7076\uff0c\u7ed9\u51fa\u4f4d\u7f6e\u3001\u7f6e\u4fe1\u5ea6\u548c\u4e25\u91cd\u5ea6\u2014\u2014\u4f46\u6ce8\u610f\uff0c\u5b83\u4e0d\u662f\u76f4\u63a5\u51fa\u62a5\u544a\u3002\u6bcf\u4e00\u5904\u6807\u6ce8\u90fd\u5fc5\u987b\u7531\u533b\u5e08\u9010\u6846\u786e\u8ba4\u6216\u9a73\u56de\uff0c\u533b\u751f\u8fd8\u53ef\u4ee5\u70b9\u51fb\u5f71\u50cf\u81ea\u5df1\u52a0\u6807\u6ce8\u3002AI \u9884\u6807\u6ce8\u662f\u5efa\u8bae\uff0c\u533b\u5e08\u590d\u6838\u624d\u662f\u88c1\u51b3\u3002\u8fd9\u6b63\u662f\u533b\u7597 AI \u5e94\u6709\u7684\u5b89\u5168\u8fb9\u754c\u3002")


def slide08(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s, LIGHT_BG)
    title_bar(s, "Demo 3 \u00b7 \u884c\u4e3a\u4fe1\u53f7\u8bc6\u522b \u2014\u2014 \u4ece\u88ab\u52a8\u7b49\u5f85\u5230\u4e3b\u52a8\u9884\u8b66",
             "\u7528\u836f / \u5c31\u533b\u6a21\u5f0f\u5206\u6790 \u00b7 \u6162\u75c5\u98ce\u9669\u5206\u7ea7", 8, accent=ORANGE)
    card(s, 0.5, 1.35, 5.6, 4.4, WHITE, line=RED, lw=1.5)
    card_title(s, 0.5, 1.35, 5.6, "\u4e3b\u52a8\u9884\u8b66\u63a8\u9001", fill=RED)
    icon(s, 0.85, 2.05, 0.7, RED, "!", size=18)
    tbox(s, 1.7, 2.05, 4.2, 0.4, [("\u6162\u75c5\u7ba1\u7406\u8bc4\u5206\u4e0b\u964d", 14, RED, True)])
    tbox(s, 1.7, 2.42, 4.2, 0.35, [("\u8fde\u7eed 3 \u6708\u8d2d\u964d\u7cd6\u836f \u00b7 \u672a\u590d\u67e5\u7cd6\u5316", 10.5, GRAY)])
    tbox(s, 0.75, 2.9, 2.0, 0.9, [("58", 40, RED, True, PP_ALIGN.CENTER)])
    tbox(s, 2.7, 2.9, 3.2, 0.9, [("\u2193 12 \u5206\n\uff08\u57fa\u7ebf 70\uff09", 12, RED, True, PP_ALIGN.CENTER)], anchor=MSO_ANCHOR.MIDDLE)
    tbox(s, 0.75, 3.95, 5.1, 0.5,
         [("\u884c\u4e3a\u4fe1\u53f7\uff1a\u7528\u836f + \u5c31\u533b + \u8d2d\u836f\u6a21\u5f0f\u7efc\u5408\u5206\u6790", 10.5, DARK, False, PP_ALIGN.CENTER)])
    lvl = rrect(s, 0.75, 4.55, 5.1, 0.95, LIGHT_BG, line=GRAY_LT, lw=1.0, radius=0.10)
    set_text(lvl.text_frame,
             [("\u5206\u7ea7\u9884\u8b66\uff1a\U0001f534 \u9ad8\u98ce\u9669    \U0001f7e1 \u4e2d\u98ce\u9669    \U0001f7e2 \u4f4e\u98ce\u9669", 12, DARK, True, PP_ALIGN.CENTER)],
             anchor=MSO_ANCHOR.MIDDLE)
    card(s, 6.35, 1.35, 6.3, 2.05, WHITE, line=GRAY_LT, lw=1.0)
    card_title(s, 6.35, 1.35, 6.3, "\u4e94\u7ef4\u5065\u5eb7\u753b\u50cf", fill=NAVY)
    dims = [("\u7528\u836f\u4f9d\u4ece", 60), ("\u590d\u67e5\u9891\u7387", 35), ("\u5c31\u533b\u89c4\u5f8b", 70), ("\u6307\u6807\u76d1\u63a7", 40), ("\u98ce\u9669\u8bc4\u5206", 58)]
    bx = 6.55; by = 2.1; bw = 1.05
    for name, val in dims:
        rect(s, bx, by, bw, 1.0, GRAY_LT)
        rect(s, bx, by + (100 - val) * 0.012, bw, val * 0.012, NAVY)
        tbox(s, bx - 0.05, 3.12, bw + 0.1, 0.22, [(name, 8.5, DARK, False, PP_ALIGN.CENTER)])
        tbox(s, bx - 0.05, by - 0.02, bw + 0.1, 0.22, [(str(val), 9, NAVY, True, PP_ALIGN.CENTER)])
        bx += 1.18
    tbox(s, 6.55, 3.28, 5.9, 0.12, [("\uff080-100\uff0c\u8d8a\u9ad8\u8d8a\u597d\uff09", 8, GRAY, False, PP_ALIGN.CENTER)])
    sugg = [("\u590d\u67e5\u63d0\u9192", "\u5efa\u8bae\u590d\u67e5\u7cd6\u5316\u8840\u7ea2\u86cb\u767d", GREEN),
            ("\u996e\u98df\u8c03\u6574", "\u4f4e\u7cd6\u996e\u98df + \u8fd0\u52a8\u5efa\u8bae", GREEN),
            ("\u7533\u8bf7\u5f85\u9047", "\u95e8\u8bca\u6162\u75c5\u5f85\u9047\u7533\u8bf7", ORANGE)]
    cw, gap = 1.95, 0.15
    x0 = 6.35
    for i, (h, b, col) in enumerate(sugg):
        x = x0 + i * (cw + gap)
        card(s, x, 3.55, cw, 2.2, WHITE, line=col, lw=1.0)
        icon(s, x + cw / 2 - 0.3, 3.7, 0.6, col, "\u2192", size=14)
        tbox(s, x + 0.1, 4.4, cw - 0.2, 0.4, [(h, 11.5, col, True, PP_ALIGN.CENTER)])
        tbox(s, x + 0.1, 4.85, cw - 0.2, 0.8, [b], size=10, color=DARK, align=PP_ALIGN.CENTER)
    notes(s, "\u7b2c\u4e09\u4e2a Demo\u3002瓯医数链 \u5206\u6790\u5f20\u5148\u751f\u7684\u7528\u836f\u884c\u4e3a\u2014\u2014\u8fde\u7eed\u4e09\u4e2a\u6708\u4e70\u964d\u7cd6\u836f\uff0c\u5374\u6ca1\u590d\u67e5\u7cd6\u5316\u8840\u7ea2\u86cb\u767d\uff0c\u6162\u75c5\u7ba1\u7406\u8bc4\u5206\u660e\u663e\u4e0b\u964d\u3002\u7cfb\u7edf\u4e0d\u662f\u7b49\u4ed6\u6765\u95ee\uff0c\u800c\u662f\u4e3b\u52a8\u628a\u9884\u8b66\u63a8\u5230\u4ed6\u9762\u524d\u3002\u4ece\u88ab\u52a8\u7b49\u7b54\u6848\u5230\u4e3b\u52a8\u6293\u4fe1\u53f7\uff0c\u8fd9\u662f\u8303\u5f0f\u8f6c\u53d8\u3002")


def slide09(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s, LIGHT_BG)
    title_bar(s, "Demo 4 \u00b7 \u60a3\u8005\u4fe1\u606f\u8fde\u63a5 \u2014\u2014 \u653f\u7b56\u7ea2\u5229\u770b\u5f97\u89c1\u3001\u62ff\u5f97\u5230",
             "\u62a5\u9500\u9884\u5ba1 \u00b7 \u653f\u7b56\u5339\u914d \u00b7 \u6743\u76ca\u95ee\u7b54 \u00b7 \u591a\u667a\u80fd\u4f53\u534f\u4f5c", 9, accent=GREEN)
    tbox(s, 0.5, 1.25, 6.0, 0.3, [("\u62a5\u9500\u9884\u5ba1 \u00b7 7 \u6b65\u8ba1\u7b97", 12, NAVY, True)])
    steps = ["\u4e0a\u4f20\u7968\u636e", "OCR \u8bc6\u522b", "\u5206\u7c7b\u6c47\u603b", "\u8d77\u4ed8\u7ebf", "\u5171\u4ed8\u6bd4\u4f8b", "\u5c01\u9876\u7ebf", "\u5e94\u62a5\u91d1\u989d"]
    n = len(steps)
    bw, gap = 1.62, 0.13
    x0 = 0.5
    for i, st in enumerate(steps):
        x = x0 + i * (bw + gap)
        col = GREEN if i == n - 1 else NAVY
        box = rrect(s, x, 1.6, bw, 0.95, col, radius=0.10)
        set_text(box.text_frame,
                 [(str(i + 1), 12, SUB_CL if col == NAVY else GREEN_LT, True, PP_ALIGN.CENTER),
                  (st, 11, WHITE, True, PP_ALIGN.CENTER)],
                 anchor=MSO_ANCHOR.MIDDLE)
        if i < n - 1:
            tbox(s, x + bw, 1.6, gap, 0.95, [("\u2192", 12, GRAY, True, PP_ALIGN.CENTER)],
                 anchor=MSO_ANCHOR.MIDDLE)
    tbox(s, 0.5, 2.62, 12.33, 0.3,
         [("\u667a\u80fd\u63d0\u9192\uff1a\u201c\u60a8\u53ef\u80fd\u7b26\u5408\u95e8\u8bca\u6162\u75c5\u5f85\u9047\u201d", 10.5, ORANGE, True, PP_ALIGN.CENTER)])
    tbox(s, 0.5, 3.0, 8.0, 0.3, [("\u653f\u7b56\u5339\u914d \u00b7 \u6309\u7701\u94b1\u91d1\u989d\u6392\u5e8f\uff08\u542b\u653f\u7b56\u539f\u6587\uff09", 12, NAVY, True)])
    policies = [("\u95e8\u8bca\u6162\u75c5\u5f85\u9047", "\u00a52,016 / \u5e74", True), ("\u5927\u75c5\u4fdd\u9669", "\u00a5800 / \u5e74", False),
               ("\u5bb6\u5ead\u8d26\u6237\u5171\u6d4e", "\u00a5350 / \u5e74", False)]
    cw, gap = 3.9, 0.3
    x0 = 0.5
    for i, (name, save, hi) in enumerate(policies):
        x = x0 + i * (cw + gap)
        col = ORANGE if hi else NAVY_LT
        card(s, x, 3.35, cw, 2.05, WHITE, line=col, lw=2.0 if hi else 1.0)
        tbox(s, x + 0.2, 3.45, cw - 0.4, 0.4, [(name, 13, NAVY, True, PP_ALIGN.CENTER)])
        tbox(s, x + 0.2, 3.9, cw - 0.4, 0.7, [(save, 22, col if hi else NAVY, True, PP_ALIGN.CENTER)])
        tag = "\u2605 \u6700\u7701\u94b1" if hi else "\u63a8\u8350"
        tbox(s, x + 0.2, 4.7, cw - 0.4, 0.55, [(tag, 10.5, col if hi else GRAY, True, PP_ALIGN.CENTER)])
    badge = rrect(s, 9.5, 3.35, 3.33, 2.05, GREEN_BG, line=GREEN, lw=1.5, radius=0.08)
    set_text(badge.text_frame,
             [("\u591a\u667a\u80fd\u4f53\u534f\u4f5c", 13, GREEN, True, PP_ALIGN.CENTER),
              ("\u5df2\u8c03\u5ea6 2 \u4e2a\u667a\u80fd\u4f53\u534f\u540c", 11, DARK, False, PP_ALIGN.CENTER),
              ("\u8111\u7535\u536b\u58eb + \u653f\u7b56\u53c2\u8c0b", 10.5, NAVY, True, PP_ALIGN.CENTER),
              ("\u201c\u8111\u7535\u5f02\u5e38\u600e\u4e48\u7701\u94b1\u201d \u2192 \u5e76\u884c\u56de\u7b54", 9.5, GRAY, False, PP_ALIGN.CENTER)],
             anchor=MSO_ANCHOR.MIDDLE)
    qa = rrect(s, 0.5, 5.6, 12.33, 0.62, NAVY, radius=0.10)
    set_text(qa.text_frame,
             [("权益问答 · 一句话查清：参保类型 / 账户余额 / 报销比例 · 主动提醒缴费年限", 13, WHITE, True, PP_ALIGN.CENTER)],
             anchor=MSO_ANCHOR.MIDDLE)
    live = rrect(s, 0.5, 6.30, 12.33, 0.58, GREEN, radius=0.10)
    set_text(live.text_frame,
             [("现场演示 · 动态新增用户 → 对话/健康画像/数字人体档案全站即时联动（< 10 秒）· 会话持久化重启可回放", 12, WHITE, True, PP_ALIGN.CENTER)],
             anchor=MSO_ANCHOR.MIDDLE)
    admin = rrect(s, 0.5, 6.98, 12.33, 0.42, NAVY, radius=0.10)
    set_text(admin.text_frame,
             [("运营管理后台 · 全用户使用统计 + 画像分析 → 精准政策传达 / 病情告知 / 营销推送", 11, WHITE, True, PP_ALIGN.CENTER)],
             anchor=MSO_ANCHOR.MIDDLE)
    notes(s, "第四个 Demo 是信息连接层。一张票据上传，报销计算七步走；政策匹配按能省多少钱排序，每一条都有政策原文；多智能体协作——当他同时问“脑电异常怎么省钱”，系统会并行调度脑电卫士和政策参谋，一次回答两个问题。动态用户管理，现场新增一个用户，不到十秒，对话、健康画像、数字人体档案全站联动；连续对话会话持久化，服务重启后历史可回放。另外我们上线了运营管理后台：全用户使用统计与画像分析，支撑精准政策传达、病情告知等服务运营。")


def slide10(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s, LIGHT_BG)
    title_bar(s, "\u56db\u5927\u6280\u672f\u521b\u65b0", "\u4e0d\u505a\u9ed1\u7bb1 AI\uff0c\u505a\u533b\u751f\u4e0e\u60a3\u8005\u80fd\u4fe1\u4efb\u7684 AI", 10, accent=GREEN)
    innov = [
        ("\u591a\u6a21\u6001\u4fe1\u53f7\u8bc6\u522b", "\u8111\u7535\u9891\u57df + \u5f71\u50cf\u75c5\u7076 + \u884c\u4e3a\u6a21\u5f0f\uff0c\u7edf\u4e00\u7531\u7f16\u6392\u667a\u80fd\u4f53\u8c03\u5ea6\uff0c\u590d\u5408\u610f\u56fe\u5e76\u884c\u6267\u884c",
         GREEN, True, "\u8111"),
        ("\u533b\u5e08\u5728\u73af Human-in-the-Loop", "AI \u5f71\u50cf\u9884\u6807\u6ce8\u5fc5\u987b\u533b\u5e08\u9010\u6846\u590d\u6838\u624d\u80fd\u51fa\u62a5\u544a\uff1b\u5efa\u8bae\u6743\u7ed9 AI\uff0c\u88c1\u51b3\u6743\u7ed9\u4eba",
         ORANGE, True, "\u533b"),
        ("\u53ef\u4fe1\u6570\u636e\u7a7a\u95f4\u5bf9\u9f50", "\u9690\u79c1\u8ba1\u7b97\u201c\u53ef\u7528\u4e0d\u53ef\u89c1\u201d + \u533a\u5757\u94fe\u5b58\u8bc1\uff0c\u5bf9\u9f50\u201c\u6570\u636e\u8981\u7d20 \u00d7 \u533b\u7597\u5065\u5eb7\u201d\u56fd\u5bb6\u6218\u7565",
         NAVY_LT, False, "\u9501"),
        ("\u5168\u94fe\u8def\u53ef\u89e3\u91ca\u6027", "\u8111\u7535\u6709\u9891\u57df\u8bc1\u636e \u00b7 \u5f71\u50cf\u6709 bbox \u5750\u6807 \u00b7 \u653f\u7b56\u6709\u539f\u6587\u5f15\u7528 \u00b7 \u62a5\u9500\u6709\u6b65\u9aa4\u63a8\u5bfc",
         NAVY_LT, False, "\u91ca"),
    ]
    cw, ch, gap = 6.0, 2.35, 0.3
    positions = [(0.5, 1.35), (6.83, 1.35), (0.5, 3.85), (6.83, 3.85)]
    for (x, y), (head, body, col, star, lab) in zip(positions, innov):
        card(s, x, y, cw, ch, WHITE, line=col, lw=2.0 if star else 1.0)
        icon(s, x + 0.25, y + 0.25, 0.7, col, lab, size=16)
        tbox(s, x + 1.1, y + 0.25, cw - 1.3, 0.5,
             [(head + (" \u2b50" if star else ""), 15, col if star else NAVY, True)])
        tbox(s, x + 0.3, y + 1.05, cw - 0.6, 1.2, [body], size=11.5, color=DARK)
    foot = rrect(s, 0.5, 6.4, 12.33, 0.7, NAVY, radius=0.10)
    set_text(foot.text_frame,
             [("\u4e0d\u505a\u9ed1\u7bb1 AI\uff0c\u505a\u533b\u751f\u4e0e\u60a3\u8005\u80fd\u4fe1\u4efb\u7684 AI", 16, WHITE, True, PP_ALIGN.CENTER)],
             anchor=MSO_ANCHOR.MIDDLE)
    notes(s, "\u56db\u5927\u521b\u65b0\uff1a\u591a\u6a21\u6001\u4fe1\u53f7\u8bc6\u522b\u8ba9\u4e09\u7c7b\u4fe1\u53f7\u7edf\u4e00\u8c03\u5ea6\uff1b\u533b\u5e08\u5728\u73af\u786e\u4fdd AI \u6c38\u8fdc\u662f\u52a9\u624b\u3001\u533b\u5e08\u6c38\u8fdc\u662f\u88c1\u5224\uff1b\u53ef\u4fe1\u6570\u636e\u7a7a\u95f4\u4fdd\u969c\u9690\u79c1\u5408\u89c4\uff1b\u5168\u94fe\u8def\u53ef\u89e3\u91ca\u8ba9\u6bcf\u4e00\u4e2a\u51b3\u7b56\u90fd\u6709\u636e\u53ef\u67e5\u3002")


def slide11(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s, LIGHT_BG)
    title_bar(s, "\u533b\u7597\u5b89\u5168 \u00b7 \u9690\u79c1\u4fdd\u62a4 \u00b7 \u4f7f\u7528\u8fb9\u754c", "\u533b\u7597\u4ea7\u54c1\u5b89\u5168\u7b2c\u4e00 \u00b7 \u660e\u786e\u4e09\u6761\u8fb9\u754c", 11, accent=ORANGE)
    cols = [
        ("\u533b\u7597\u5b89\u5168", NAVY, "\u76fe",
         ["\u5f71\u50cf\u62a5\u544a\u5fc5\u987b\u533b\u5e08\u590d\u6838\u7b7e\u5b57", "AI \u9884\u6807\u6ce8\u4ec5\u4f5c\u8f85\u52a9\u5efa\u8bae", "\u8111\u7535/\u5065\u5eb7\u9884\u8b66\u4e3a\u201c\u63d0\u793a\u201d\u975e\u201c\u8bca\u65ad\u201d", "\u660e\u786e\u514d\u8d23\u8fb9\u754c\u58f0\u660e"]),
        ("\u9690\u79c1\u4fdd\u62a4", GREEN, "\u9501",
         ["\u6570\u636e\u6388\u6743\u7ba1\u7406\uff1a\u7528\u6237\u81ea\u4e3b\u51b3\u5b9a", "\u9690\u79c1\u8ba1\u7b97\u201c\u53ef\u7528\u4e0d\u53ef\u89c1\u201d", "\u533a\u5757\u94fe\u5b58\u8bc1\u5168\u94fe\u8def\u8ffd\u6eaf", "\u5171\u4eab\u8303\u56f4\u53ef\u63a7\u53ef\u64a4\u56de"]),
        ("\u4f7f\u7528\u8fb9\u754c", ORANGE, "\u754c",
         ["\u4e0d\u652f\u6301\u5904\u65b9\u3001\u786e\u8bca\u7b49\u66ff\u4ee3\u533b\u7597", "\u5371\u6025\u4fe1\u53f7\u5f15\u5bfc\u7ebf\u4e0b\u5c31\u533b", "\u4ea7\u54c1\u5185\u7f6e\u8fb9\u754c\u58f0\u660e", "\u4e0d\u505a\u51b3\u7b56\u53ea\u505a\u63d0\u793a"]),
    ]
    cw, gap = 3.9, 0.3
    x0 = 0.5
    for i, (head, col, lab, pts) in enumerate(cols):
        x = x0 + i * (cw + gap)
        card(s, x, 1.4, cw, 4.45, WHITE, line=col, lw=1.5)
        rect(s, x, 1.4, cw, 0.7, col)
        icon(s, x + 0.25, 1.55, 0.4, WHITE, lab, size=12)
        tbox(s, x + 0.75, 1.4, cw - 0.9, 0.7, [(head, 16, WHITE, True)], anchor=MSO_ANCHOR.MIDDLE)
        lines = [("\u2022 " + p, 11, DARK) for p in pts]
        tbox(s, x + 0.25, 2.3, cw - 0.5, 3.4, lines)
    decl = rrect(s, 0.5, 6.05, 12.33, 0.7, LIGHT_BG, line=GRAY, lw=1.0, radius=0.10)
    set_text(decl.text_frame,
             [("\u4ea7\u54c1\u4f7f\u7528\u8fb9\u754c\u58f0\u660e\uff1a\u672c\u4ea7\u54c1\u6240\u6709\u8f93\u51fa\u5747\u4e3a\u5065\u5eb7\u63d0\u793a\uff0c\u4e0d\u6784\u6210\u533b\u7597\u8bca\u65ad\u6216\u5904\u65b9\u5efa\u8bae\uff0c\u5371\u6025\u60c5\u51b5\u8bf7\u7acb\u5373\u5c31\u533b", 10.5, GRAY, False, PP_ALIGN.CENTER)],
             anchor=MSO_ANCHOR.MIDDLE)
    notes(s, "\u533b\u7597\u4ea7\u54c1\u5b89\u5168\u7b2c\u4e00\u3002\u6211\u4eec\u660e\u786e\u4e09\u6761\u8fb9\u754c\uff1a\u5f71\u50cf\u62a5\u544a\u5fc5\u987b\u533b\u5e08\u590d\u6838\u7b7e\u5b57\uff1b\u6240\u6709\u9884\u8b66\u90fd\u662f\u63d0\u793a\u4e0d\u662f\u8bca\u65ad\uff1b瓯医数链 \u4e0d\u505a\u5904\u65b9\u3001\u4e0d\u505a\u786e\u8bca\uff0c\u5371\u6025\u4fe1\u53f7\u4e00\u5f8b\u5f15\u5bfc\u7ebf\u4e0b\u5c31\u533b\u3002\u6570\u636e\u65b9\u9762\uff0c\u7528\u6237\u53ef\u4ee5\u81ea\u4e3b\u6388\u6743\uff0c\u5168\u7a0b\u5b58\u8bc1\u53ef\u8ffd\u6eaf\u3002")


def slide12(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s, LIGHT_BG)
    title_bar(s, "\u4ece\u201c\u7b49\u751f\u75c5\u201d\u5230\u201c\u65e9\u53d1\u73b0\u201d", "\u8303\u5f0f\u8f6c\u53d8\uff1a\u4e3b\u52a8\u6293\u4fe1\u53f7\uff0c\u63d0\u524d\u6355\u6349\u5f02\u5e38", 12, accent=GREEN)
    card(s, 0.5, 1.4, 5.8, 3.1, WHITE, line=GRAY_LT, lw=1.0)
    card_title(s, 0.5, 1.4, 5.8, "\u8303\u5f0f\u5bf9\u6bd4", fill=NAVY)
    old = rrect(s, 0.8, 2.0, 5.2, 1.0, GRAY_LT, radius=0.10)
    set_text(old.text_frame,
             [("\u88ab\u52a8 \u00b7 \u7b49\u751f\u75c5", 13, GRAY, True, PP_ALIGN.CENTER),
              ("\u65e0\u75c7\u72b6\u4e0d\u68c0\u67e5 \u2192 \u51fa\u73b0\u75c7\u72b6\u624d\u5c31\u533b \u2192 \u6210\u672c\u9ad8\u3001\u53d1\u73b0\u665a", 10, GRAY, False, PP_ALIGN.CENTER)],
             anchor=MSO_ANCHOR.MIDDLE)
    new = rrect(s, 0.8, 3.15, 5.2, 1.2, GREEN, radius=0.10)
    set_text(new.text_frame,
             [("\u4e3b\u52a8 \u00b7 \u65e9\u53d1\u73b0", 13, WHITE, True, PP_ALIGN.CENTER),
              ("\u8111\u7535/\u5f71\u50cf/\u884c\u4e3a\u4fe1\u53f7\u63d0\u524d\u6355\u6349\u5f02\u5e38 \u2192 \u964d\u4f4e\u53d1\u73b0\u6210\u672c", 10, WHITE, False, PP_ALIGN.CENTER)],
             anchor=MSO_ANCHOR.MIDDLE)
    vals = [("\u65e9\u53d1\u73b0\u65e9\u5e72\u9884", "\u4e09\u7c7b\u4fe1\u53f7\u63d0\u524d\u6355\u6349\u5f02\u5e38", GREEN),
            ("\u57fa\u5c42\u8d4b\u80fd", "AI \u9884\u6807\u6ce8\u7f13\u89e3\u9605\u7247\u538b\u529b", NAVY_LT),
            ("\u4fe1\u606f\u666e\u60e0", "\u7ffb\u8bd1\u6210\u60a3\u8005\u770b\u5f97\u61c2\u7684\u201c\u4eba\u8bdd\u201d", ORANGE),
            ("\u6570\u636e\u4ef7\u503c\u91ca\u653e", "\u9690\u79c1\u8ba1\u7b97\u4e0b\u91ca\u653e\u6570\u636e\u8981\u7d20", NAVY)]
    cw, ch, gap = 3.0, 1.4, 0.2
    x0, y0 = 6.55, 1.4
    for i, (h, b, col) in enumerate(vals):
        x = x0 + (i % 2) * (cw + gap)
        y = y0 + (i // 2) * (ch + gap)
        card(s, x, y, cw, ch, WHITE, line=col, lw=1.5)
        icon(s, x + 0.2, y + 0.2, 0.5, col, "\u2605", size=12)
        tbox(s, x + 0.8, y + 0.18, cw - 0.95, 0.5, [(h, 13, NAVY, True)])
        tbox(s, x + 0.25, y + 0.78, cw - 0.5, 0.55, [b], size=10.5, color=DARK)
    kws = [("\u65e9\u53d1\u73b0", GREEN), ("\u5f3a\u57fa\u5c42", NAVY), ("\u60e0\u6c11\u751f", ORANGE)]
    kw_y = 4.65
    cw2, gap2 = 3.9, 0.3
    x0 = 0.5
    for i, (kw, col) in enumerate(kws):
        x = x0 + i * (cw2 + gap2)
        box = rrect(s, x, kw_y, cw2, 0.85, col, radius=0.10)
        set_text(box.text_frame, [(kw, 22, WHITE, True, PP_ALIGN.CENTER)], anchor=MSO_ANCHOR.MIDDLE)
    tags = ["\u6570\u636e\u8981\u7d20", "\u5065\u5eb7\u4e2d\u56fd", "\u5206\u7ea7\u8bca\u7597", "\u53ef\u4fe1\u6570\u636e\u7a7a\u95f4", "\u533b\u7597 AI"]
    tx = 0.5
    ty = 5.75
    for tag in tags:
        w = 0.3 + len(tag) * 0.18
        tg = rrect(s, tx, ty, w, 0.45, LIGHT_BG, line=NAVY_LT, lw=1.0, radius=0.30)
        set_text(tg.text_frame, [(tag, 11, NAVY, True, PP_ALIGN.CENTER)], anchor=MSO_ANCHOR.MIDDLE)
        tx += w + 0.2
    notes(s, "瓯医数链 \u7684\u4ef7\u503c\u662f\u8ba9\u5065\u5eb7\u7ba1\u7406\u4ece\u201c\u7b49\u751f\u75c5\u201d\u53d8\u6210\u201c\u65e9\u53d1\u73b0\u201d\uff1a\u57fa\u5c42\u533b\u751f\u6709\u4e86 AI \u52a9\u624b\uff0c\u60a3\u8005\u6709\u4e86\u80fd\u770b\u61c2\u4fe1\u53f7\u3001\u62ff\u5f97\u5230\u7ea2\u5229\u7684\u7ba1\u5bb6\u3002\u8fd9\u662f\u5bf9\u201c\u5065\u5eb7\u4e2d\u56fd\u201d\u6700\u76f4\u63a5\u7684\u54cd\u5e94\u3002")


def slide13(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s, LIGHT_BG)
    title_bar(s, "\u4e09\u5c42\u5546\u4e1a\u6a21\u5f0f", "\u4fe1\u53f7\u8bc6\u522b\u662f\u80fd\u529b\u5e95\u5ea7 \u00b7 \u6570\u636e\u670d\u52a1\u662f\u589e\u957f\u5f15\u64ce", 13, accent=GREEN)
    tiers = [
        ("To C \u00b7 \u4e2a\u4eba / \u5bb6\u5ead", "\u57fa\u7840\u4fe1\u53f7\u8bc6\u522b\u514d\u8d39 \u00b7 \u5065\u5eb7\u7ba1\u7406\u589e\u503c\u4ed8\u8d39",
         "\u4f1a\u5458\u8ba2\u9605 / \u5355\u6b21\u670d\u52a1", NAVY, 9.5, 1.75),
        ("To B \u00b7 \u4f53\u68c0 / \u5546\u4fdd / \u836f\u4f01", "\u591a\u6a21\u6001\u4fe1\u53f7\u8bc6\u522b\u80fd\u529b\u8f93\u51fa \u00b7 \u771f\u5b9e\u4e16\u754c\u6570\u636e\u5206\u6790",
         "\u6570\u636e\u670d\u52a1\u8d39 / API \u8c03\u7528\uff08\u8131\u654f\u5408\u89c4\uff09", GREEN, 10.5, 2.85),
        ("To G \u00b7 \u533b\u9662 / \u536b\u5065 / \u533b\u4fdd", "\u5f71\u50cf AI \u8f85\u52a9\u9605\u7247 \u00b7 \u667a\u80fd\u5ba2\u670d\u66ff\u4ee3\u4eba\u5de5\u7a97\u53e3",
         "\u6309\u8c03\u7528\u91cf / \u5e74\u670d\u52a1\u8d39", ORANGE, 11.5, 3.95),
    ]
    for (head, body, fee, col, w, y) in tiers:
        x = (SW - w) / 2
        box = rrect(s, x, y, w, 1.05, col, radius=0.10)
        set_text(box.text_frame,
                 [(head, 15, WHITE, True, PP_ALIGN.CENTER),
                  (body, 10, WHITE, False, PP_ALIGN.CENTER),
                  (fee, 10, WHITE if col == GREEN else GREEN_LT, True, PP_ALIGN.CENTER)],
                 anchor=MSO_ANCHOR.MIDDLE)
    note = rrect(s, 0.5, 5.15, 12.33, 0.7, NAVY, radius=0.10)
    set_text(note.text_frame,
             [("\u6536\u5165\u7ed3\u6784\u9884\u4f30\uff1aTo G \u80fd\u529b\u8f93\u51fa\u4e3a\u4e3b \u00b7 To B \u6570\u636e\u670d\u52a1\u4e3a\u589e\u957f\u5f15\u64ce \u00b7 To C \u589e\u503c\u8ba2\u9605\u4e3a\u957f\u5c3e", 13, WHITE, True, PP_ALIGN.CENTER)],
             anchor=MSO_ANCHOR.MIDDLE)
    pts = [("\u80fd\u529b\u5e95\u5ea7", "\u591a\u6a21\u6001\u4fe1\u53f7\u8bc6\u522b\u7edf\u4e00\u8f93\u51fa", GREEN),
           ("\u589e\u957f\u5f15\u64ce", "\u8131\u654f\u6570\u636e\u670d\u52a1 + API \u8c03\u7528", ORANGE),
           ("\u666e\u60e0\u5165\u53e3", "To C \u514d\u8d39\u5f15\u6d41 + \u589e\u503c\u53d8\u73b0", NAVY)]
    cw, gap = 3.9, 0.3
    x0 = 0.5
    for i, (h, b, col) in enumerate(pts):
        x = x0 + i * (cw + gap)
        card(s, x, 6.0, cw, 1.35, WHITE, line=col, lw=1.0)
        tbox(s, x + 0.25, 6.08, cw - 0.5, 0.35, [(h, 12, col, True, PP_ALIGN.CENTER)])
        tbox(s, x + 0.25, 6.45, cw - 0.5, 0.85, [b], size=10, color=DARK, align=PP_ALIGN.CENTER)
    notes(s, "\u5546\u4e1a\u5316\u5206\u4e09\u5c42\uff1a\u9762\u5411\u533b\u9662\u548c\u533b\u4fdd\u673a\u6784\u8f93\u51fa\u8bc6\u522b\u80fd\u529b\uff1b\u9762\u5411\u4f53\u68c0\u3001\u5546\u4fdd\u3001\u836f\u4f01\u63d0\u4f9b\u8131\u654f\u6570\u636e\u5206\u6790\uff1b\u9762\u5411\u4e2a\u4eba\u57fa\u7840\u529f\u80fd\u514d\u8d39\u3001\u589e\u503c\u670d\u52a1\u4ed8\u8d39\u3002\u4fe1\u53f7\u8bc6\u522b\u662f\u80fd\u529b\u5e95\u5ea7\uff0c\u6570\u636e\u670d\u52a1\u662f\u589e\u957f\u5f15\u64ce\u3002")


def slide14(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s, LIGHT_BG)
    title_bar(s, "\u4ece\u9ed1\u5ba2\u677e\u5230\u771f\u5b9e\u4e34\u5e8a\u573a\u666f", "\u843d\u5730\u4e09\u6b65\u8d70 \u00b7 \u5173\u952e\u4f9d\u8d56\uff1a\u771f\u5b9e\u6570\u636e + \u4e34\u5e8a\u5408\u4f5c + \u5408\u89c4\u8def\u5f84", 14, accent=GREEN)
    rect(s, 0.8, 3.35, 11.7, 0.08, GRAY_LT)
    phases = [
        ("Phase 1", "1\u20133 \u4e2a\u6708", "\u5408\u4f5c\u8bd5\u70b9", GREEN, True,
         ["\u4e0e\u57fa\u5c42\u533b\u7597/\u4f53\u68c0\u673a\u6784\u8054\u5408\u8bd5\u70b9", "\u5f71\u50cf\u590d\u6838\u5de5\u4f5c\u6d41\u9a8c\u8bc1", "\u63a5\u5165\u771f\u5b9e\u8131\u654f\u6570\u636e"]),
        ("Phase 2", "3\u20136 \u4e2a\u6708", "\u591a\u4e2d\u5fc3\u6269\u5c55", NAVY, False,
         ["\u8986\u76d6 2\u20133 \u5bb6\u533b\u7597\u673a\u6784", "\u6269\u5c55\u5f71\u50cf\u7c7b\u578b\u4e0e\u8111\u7535\u8bbe\u5907", "\u6536\u96c6\u533b\u751f/\u60a3\u8005\u53cd\u9988\u8fed\u4ee3"]),
        ("Phase 3", "6\u201312 \u4e2a\u6708", "\u4ea7\u54c1\u5316", ORANGE, False,
         ["\u533b\u7597\u5668\u68b0\u8f6f\u4ef6\uff08SaMD\uff09\u6ce8\u518c", "\u53ef\u4fe1\u6570\u636e\u7a7a\u95f4\u89c4\u6a21\u5316\u90e8\u7f72", "\u5546\u4e1a\u5316\u843d\u5730"]),
    ]
    cw, gap = 3.9, 0.3
    x0 = 0.5
    for i, (ph, dur, sub, col, here, pts) in enumerate(phases):
        x = x0 + i * (cw + gap)
        oval(s, x + cw / 2 - 0.25, 3.15, 0.5, col)
        if here:
            oval(s, x + cw / 2 - 0.35, 3.05, 0.7, col, line=col, lw=2.0)
        card(s, x, 3.75, cw, 2.6, WHITE, line=col, lw=1.5)
        rect(s, x, 3.75, cw, 0.65, col)
        tbox(s, x + 0.2, 3.77, cw - 0.4, 0.6,
             [(ph + " \u00b7 " + dur, 13, WHITE, True, PP_ALIGN.CENTER)], anchor=MSO_ANCHOR.MIDDLE)
        tbox(s, x + 0.25, 4.5, cw - 0.5, 0.35, [(sub, 13, col, True, PP_ALIGN.CENTER)])
        lines = [("\u2022 " + p, 10.5, DARK) for p in pts]
        tbox(s, x + 0.25, 4.9, cw - 0.5, 1.4, lines)
        if here:
            tag = rrect(s, x + cw / 2 - 0.7, 2.7, 1.4, 0.38, ORANGE, radius=0.30)
            set_text(tag.text_frame, [("\u6211\u4eec\u5728\u8fd9\u91cc", 9, WHITE, True, PP_ALIGN.CENTER)], anchor=MSO_ANCHOR.MIDDLE)
    dep = rrect(s, 0.5, 6.5, 12.33, 0.65, NAVY, radius=0.10)
    set_text(dep.text_frame,
             [("\u5173\u952e\u4f9d\u8d56\uff1a\u771f\u5b9e\u8131\u654f\u6570\u636e \u00b7 \u4e34\u5e8a\u673a\u6784\u5408\u4f5c \u00b7 SaMD \u5408\u89c4\u6ce8\u518c\u8def\u5f84 \u2014\u2014 \u7b97\u6cd5\u4e0e\u5de5\u4f5c\u6d41\u5df2\u5728\u9ed1\u5ba2\u677e\u7248\u672c\u5b8c\u6574\u8dd1\u901a", 11.5, WHITE, True, PP_ALIGN.CENTER)],
             anchor=MSO_ANCHOR.MIDDLE)
    notes(s, "\u843d\u5730\u5206\u4e09\u6b65\uff1a\u5148\u4e0e\u57fa\u5c42\u533b\u7597\u6216\u4f53\u68c0\u673a\u6784\u8bd5\u70b9\u5f71\u50cf\u590d\u6838\uff0c\u518d\u6269\u5c55\u5230\u591a\u4e2d\u5fc3\u548c\u66f4\u591a\u4fe1\u53f7\u7c7b\u578b\uff0c\u6700\u540e\u8d70\u5411\u533b\u7597\u5668\u68b0\u8f6f\u4ef6\u7684\u4ea7\u54c1\u5316\u3002\u5173\u952e\u4f9d\u8d56\u662f\u771f\u5b9e\u6570\u636e\u548c\u4e34\u5e8a\u5408\u4f5c\uff0c\u800c\u6211\u4eec\u7684\u7b97\u6cd5\u4e0e\u5de5\u4f5c\u6d41\u5df2\u7ecf\u5728\u9ed1\u5ba2\u677e\u7248\u672c\u91cc\u5b8c\u6574\u8dd1\u901a\u3002")


def slide15(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s, LIGHT_BG)
    title_bar(s, "\u56e2\u961f", "瓯医数链 Team \u00b7 \u591a\u6a21\u6001\u533b\u7597\u4fe1\u53f7\u8bc6\u522b\u56e2\u961f", 15, accent=GREEN)
    oval(s, SW / 2 - 1.1, 1.7, 2.2, NAVY)
    tbox(s, SW / 2 - 1.1, 1.7, 2.2, 2.2, [("MS", 40, WHITE, True, PP_ALIGN.CENTER)],
         anchor=MSO_ANCHOR.MIDDLE)
    tbox(s, 1.0, 4.05, 11.3, 0.45, [("瓯医数链 Team", 22, NAVY, True, PP_ALIGN.CENTER)])
    tbox(s, 1.0, 4.5, 11.3, 0.4, [("\u5168\u6808\u5f00\u53d1 + \u4ea7\u54c1\u8bbe\u8ba1 + AI \u5e94\u7528\u5f00\u53d1", 13, GRAY, False, PP_ALIGN.CENTER)])
    tbox(s, 1.0, 5.05, 11.3, 0.3, [("\u6280\u672f\u6808", 12, NAVY, True, PP_ALIGN.CENTER)])
    tags = ["Python", "FastAPI", "Next.js", "React", "Multi-Agent", "EEG \u4fe1\u53f7\u5904\u7406", "\u533b\u5b66\u5f71\u50cf\u5de5\u4f5c\u6d41", "LLM \u5e94\u7528", "Prompt \u5de5\u7a0b"]
    tx0 = 0.8
    ty = 5.4
    tx = tx0
    row_max = 12.0
    for tag in tags:
        w = 0.45 + len(tag) * 0.13
        if tx + w > row_max:
            tx = tx0
            ty += 0.6
        col = GREEN if tag in ("Multi-Agent", "EEG \u4fe1\u53f7\u5904\u7406", "\u533b\u5b66\u5f71\u50cf\u5de5\u4f5c\u6d41") else NAVY_LT
        tg = rrect(s, tx, ty, w, 0.45, WHITE, line=col, lw=1.0, radius=0.30)
        set_text(tg.text_frame, [(tag, 11, col, True, PP_ALIGN.CENTER)], anchor=MSO_ANCHOR.MIDDLE)
        tx += w + 0.2
    foot = rrect(s, 0.5, 6.6, 12.33, 0.65, GREEN, radius=0.10)
    set_text(foot.text_frame,
             [("\u4ece 0 \u5230 1\uff0c\u628a\u60f3\u6cd5\u53d8\u6210\u53ef\u8fd0\u884c\u7684\u533b\u7597\u4ea7\u54c1", 15, WHITE, True, PP_ALIGN.CENTER)],
             anchor=MSO_ANCHOR.MIDDLE)
    notes(s, "\u6211\u4eec\u56e2\u961f\u736b\u7acb\u5b8c\u6210\u4e86\u67b6\u6784\u3001\u524d\u540e\u7aef\u3001EEG \u4fe1\u53f7\u5904\u7406\u4e0e\u533b\u5b66\u5f71\u50cf\u5de5\u4f5c\u6d41\u7684\u5168\u90e8\u5f00\u53d1\uff0c\u4ece 0 \u5230 1 \u628a\u5b83\u53d8\u6210\u4e86\u53ef\u8fd0\u884c\u7684\u4ea7\u54c1\u3002")


def slide16(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s, NAVY)
    rect(s, 0, 0, SW, 0.22, GREEN)
    rect(s, 0, SH - 0.22, SW, 0.22, NAVY_DARK)
    tbox(s, 1.0, 1.0, 11.3, 0.6, [("瓯医数链", 30, WHITE, True, PP_ALIGN.CENTER)])
    q = rrect(s, 1.5, 2.2, 10.33, 2.0, NAVY_DARK, line=GREEN, lw=1.5, radius=0.08)
    set_text(q.text_frame,
             [("\u8ba9\u5173\u952e\u533b\u7597\u4fe1\u53f7", 34, WHITE, True, PP_ALIGN.CENTER),
              ("\u4e0d\u518d\u88ab\u9519\u8fc7", 38, ORANGE, True, PP_ALIGN.CENTER)],
             anchor=MSO_ANCHOR.MIDDLE)
    sg = rrect(s, 2.5, 4.5, 8.33, 0.7, GREEN, radius=0.10)
    set_text(sg.text_frame,
             [("\u8bc6\u522b\u4fe1\u53f7 \u00b7 \u5b88\u62a4\u5065\u5eb7 \u00b7 \u8fde\u63a5\u8d44\u6e90", 18, WHITE, True, PP_ALIGN.CENTER)],
             anchor=MSO_ANCHOR.MIDDLE)
    sigs = [("\u8111\u6ce2", GREEN), ("\u5f71\u50cf", ORANGE), ("\u5fc3\u8df3", GREEN)]
    cw, gap = 1.3, 0.4
    total = cw * 3 + gap * 2
    x0 = (SW - total) / 2
    for i, (lab, col) in enumerate(sigs):
        icon(s, x0 + i * (cw + gap), 5.45, cw, col, lab, size=14)
    agents = [("\u8111\u7535", GREEN), ("\u5f71\u50cf", GREEN), ("\u5065\u5eb7", GREEN),
              ("\u6743\u76ca", NAVY_LT), ("\u62a5\u9500", NAVY_LT), ("\u653f\u7b56", NAVY_LT), ("\u5b89\u5168", ORANGE)]
    n = len(agents)
    aw, g2 = 0.9, 0.2
    tot = aw * n + g2 * (n - 1)
    st = (SW - tot) / 2
    for i, (lab, col) in enumerate(agents):
        icon(s, st + i * (aw + g2), 6.35, aw, col, lab, size=9)
    tbox(s, 0.5, 7.0, 12.3, 0.3, [("\u8c22\u8c22\u5927\u5bb6\uff01", 13, SUB_CL, True, PP_ALIGN.CENTER)])
    notes(s, "\u56de\u5230\u5f20\u5148\u751f\u7684\u6545\u4e8b\u3002\u6709\u4e86 瓯医数链\uff0c\u538b\u529b\u6709\u4e86\u6570\u5b57\uff0c\u75c5\u7076\u6709\u4eba\u628a\u5173\uff0c\u653f\u7b56\u7ea2\u5229\u4e0d\u518d\u9519\u8fc7\u3002\u8ba9\u5173\u952e\u533b\u7597\u4fe1\u53f7\uff0c\u4e0d\u518d\u88ab\u9519\u8fc7\u3002\u8c22\u8c22\u5927\u5bb6\uff01")


def main():
    prs = Presentation()
    prs.slide_width = Inches(SW)
    prs.slide_height = Inches(SH)
    builders = [slide01, slide02, slide03, slide04, slide05, slide06, slide07,
                slide08, slide09, slide10, slide11, slide12, slide13, slide14,
                slide15, slide16]
    for b in builders:
        b(prs)
    out = r"d:\APPs\VentureDhealthcare\docs\瓯医数链.pptx"
    prs.save(out)
    size = os.path.getsize(out)
    print("OK -> " + out)
    print("slides: " + str(len(prs.slides)))
    print("size: " + str(size) + " bytes (" + str(round(size / 1024, 1)) + " KB)")


if __name__ == "__main__":
    main()
