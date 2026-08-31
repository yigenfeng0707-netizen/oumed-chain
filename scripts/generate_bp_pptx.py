"""瓯医数链初赛路演 PPT 生成（python-pptx，16:9）

基于 docs/BP-瓯医数链-初赛版.md 的内容与 docs/screenshots/ 实测截图。
运行：backend/.venv/Scripts/python scripts/generate_bp_pptx.py
输出：docs/瓯医数链-初赛路演.pptx
"""

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Emu, Inches, Pt

ROOT = Path(__file__).resolve().parent.parent
SHOTS = ROOT / "docs" / "screenshots"
OUT = ROOT / "docs" / "瓯医数链-初赛路演.pptx"

CYAN = RGBColor(0x08, 0x91, 0xB2)
DARK = RGBColor(0x1E, 0x29, 0x3B)
GRAY = RGBColor(0x64, 0x74, 0x8B)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]


def add_slide():
    return prs.slides.add_slide(BLANK)


def rect(slide, x, y, w, h, color):
    from pptx.enum.shapes import MSO_SHAPE

    sh = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    sh.fill.solid()
    sh.fill.fore_color.rgb = color
    sh.line.fill.background()
    return sh


def text(slide, x, y, w, h, content, size=18, bold=False, color=DARK,
         align=PP_ALIGN.LEFT, line_spacing=1.15):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = True
    lines = content.split("\n")
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = line_spacing
        r = p.add_run()
        r.text = line
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.color.rgb = color
    return box


def bullets(slide, x, y, w, h, items, size=15, color=DARK, gap=6):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = True
    for i, it in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(gap)
        r = p.add_run()
        r.text = "• " + it
        r.font.size = Pt(size)
        r.font.color.rgb = color
    return box


def header(slide, title, subtitle=""):
    rect(slide, 0, 0, 13.333, 0.12, CYAN)
    text(slide, 0.55, 0.32, 12.2, 0.7, title, size=28, bold=True)
    if subtitle:
        text(slide, 0.57, 0.98, 12.2, 0.5, subtitle, size=13, color=GRAY)


def page_no(slide, n):
    text(slide, 12.6, 7.05, 0.6, 0.35, str(n), size=11, color=GRAY, align=PP_ALIGN.RIGHT)


# ---------- 1 封面 ----------
s = add_slide()
rect(s, 0, 0, 13.333, 7.5, DARK)
rect(s, 0, 5.1, 13.333, 0.06, CYAN)
text(s, 1.0, 1.7, 11.3, 1.2, "瓯医数链 OuMedTrust", size=54, bold=True, color=WHITE)
text(s, 1.0, 3.0, 11.3, 0.8, "医疗数据要素可信流通平台", size=26, color=CYAN)
text(s, 1.0, 3.9, 11.3, 0.6, "让医疗数据「可用不可见、可控可计量」", size=18, color=RGBColor(0xCB, 0xD5, 0xE1))
text(s, 1.0, 5.5, 11.3, 0.6, "第二届全球技术创新大赛 · AI+医疗专题赛（赛道二：医疗大模型与数据）",
     size=15, color=RGBColor(0x94, 0xA3, 0xB8))
text(s, 1.0, 6.1, 11.3, 0.6, "2026 · 中国温州", size=13, color=RGBColor(0x94, 0xA3, 0xB8))

# ---------- 2 痛点 ----------
s = add_slide()
header(s, "三重困境：医疗数据动不了、不值钱、没基建")
cols = [
    ("不敢共享", "数据安全法/个保法硬约束\n传统数据大池化模式\n在医院侧已经走不通"),
    ("不会增值", "基层医院数据缺失率 12%-32%\n单院建模 AUC 仅 0.69\n数据资产沉睡、无法变现"),
    ("没有基建", "定价、授权、分成、监管\n缺乏可用不可见的技术底座\n与合规流通流程"),
]
for i, (t, body) in enumerate(cols):
    x = 0.7 + i * 4.1
    rect(s, x, 1.9, 3.8, 0.9, CYAN)
    text(s, x + 0.2, 2.05, 3.4, 0.6, t, size=22, bold=True, color=WHITE)
    rect(s, x, 2.8, 3.8, 2.6, RGBColor(0xF1, 0xF5, 0xF9))
    text(s, x + 0.25, 3.0, 3.3, 2.2, body, size=15, color=DARK)
text(s, 0.7, 5.9, 12.0, 1.0,
     "政策机遇：数据要素×医疗健康为国家行动方向，浙江省为市场化配置改革试点\n"
     "大赛主办方（温州市经信局）核心 KPI 即数据要素；中国（温州）智能谷是天然落地载体",
     size=15, bold=False, color=CYAN)
page_no(s, 2)

# ---------- 3 方案架构 ----------
s = add_slide()
header(s, "解决方案：1 个可信数据底座 + N 个医疗 AI 应用生态")
layers = [
    ("应用层", "健康卫士 · 影像卫士 · 档案管家 · 政策参谋 等 8 个医疗智能体（数据消费方与民生出口）", RGBColor(0x0E, 0xA5, 0xE9)),
    ("流通层", "数据产品目录 → 用途限定授权 → 交易结算 → 收益分成（医院70% / 平台20% / 贡献者10%）", RGBColor(0x08, 0x91, 0xB2)),
    ("协作层", "联邦学习（数据不出院联合建模）+ AI 病历治理 Copilot（本地大模型脱敏与结构化）", RGBColor(0x1D, 0x4E, 0xD8)),
    ("合规层", "差分隐私 + 审计存证链（sha256 串联防篡改）+ 监管方实时看板（隐私事件监控）", RGBColor(0x0F, 0x17, 0x2A)),
]
for i, (name, body, color) in enumerate(layers):
    y = 1.75 + i * 1.28
    rect(s, 0.7, y, 2.1, 1.05, color)
    text(s, 0.9, y + 0.25, 1.8, 0.6, name, size=20, bold=True, color=WHITE)
    rect(s, 2.8, y, 9.8, 1.05, RGBColor(0xF8, 0xFA, 0xFC))
    text(s, 3.05, y + 0.22, 9.4, 0.7, body, size=15)
text(s, 0.7, 6.9, 12.0, 0.5,
     "闭环：社区卫生中心病历 → AI治理脱敏 → 上架流通 → 联邦建模 → 模型被应用采购 → 全程存证 → 监管可见",
     size=14, bold=True, color=CYAN)
page_no(s, 3)

# ---------- 4 联邦学习（截图） ----------
s = add_slide()
header(s, "联邦学习：数据不出院，追平集中训练上界", "三家异构医院（三甲/县医院/社区中心）联合训练心衰30天再入院风险模型")
if (SHOTS / "01-联邦协作网络-训练结果.png").exists():
    s.shapes.add_picture(str(SHOTS / "01-联邦协作网络-训练结果.png"), Inches(0.55), Inches(1.6), height=5.0)
bullets(s, 8.3, 1.9, 4.6, 4.8, [
    "全局 AUC 0.7018，超过任何单院模型（0.6896-0.7013）",
    "追平「数据大池化」集中训练上界 0.7012——而现实中集中训练不可行",
    "UCI 心脏病 297 例真实患者复验：AUC 0.9091 追平集中上界 0.9019",
    "逐院公平性：三家医院全部获益，基层医院提升最显著",
    "差分隐私分级可调：轻噪声档仅损失 0.01 AUC",
    "94 毫秒完成一次联邦任务，纯 CPU、无 GPU 依赖",
], size=15)
page_no(s, 4)

# ---------- 5 AI 病历治理（截图） ----------
s = add_slide()
header(s, "AI 病历治理 Copilot：院内网完成，数据不出院", "PHI 规则脱敏（零漏检）+ 本地大模型结构化（qwen3:4b）+ 规则引擎自动兜底")
if (SHOTS / "02-AI病历治理-脱敏对比.png").exists():
    s.shapes.add_picture(str(SHOTS / "02-AI病历治理-脱敏对比.png"), Inches(0.55), Inches(1.6), height=5.0)
bullets(s, 8.3, 1.9, 4.6, 4.8, [
    "四类 PHI 实体（身份证/手机号/姓名/住院号）确定性掩码，保留临床语义",
    "非结构化病历 → 7 字段标准化 JSON：患者/主诉/诊断/体征/用药/既往史",
    "本地 qwen3:4b 院内网推理约 10-60 秒，断网可跑、云端密钥零依赖",
    "LLM 失效自动降级规则引擎——治理系统永不宕机",
    "治理产物 = 可流通数据原料，直接对接数据要素市场",
], size=15)
page_no(s, 5)

# ---------- 6 数据要素市场（截图） ----------
s = add_slide()
header(s, "数据要素市场：目录 → 授权 → 分成 → 监管", "收益分成：医院 70% / 平台 20% / 数据贡献者 10%")
if (SHOTS / "03-数据要素市场-监管看板.png").exists():
    s.shapes.add_picture(str(SHOTS / "03-数据要素市场-监管看板.png"), Inches(0.55), Inches(1.6), height=5.0)
bullets(s, 8.3, 1.9, 4.6, 4.8, [
    "5 类在售产品：数据集 / 治理产物 / 模型API / 算法服务",
    "用途限定授权：买方、用途、金额全流程留痕",
    "每笔交易 sha256 存证上链，篡改即断链，监管一键校验",
    "监管方看板：流通金额、收益分配、隐私事件（0 起）实时可见",
    "我们自己的联邦模型就是市场商品——生态自造血",
], size=15)
page_no(s, 6)

# ---------- 7 实验数据（真实数据复验优先，合成基准并列） ----------
s = add_slide()
header(s, "实验验证：真实公开数据复验 + 合成基准，均可复现")
rows = [
    ("方案（UCI 心脏病 297 例真实患者）", "全局测试 AUC"),
    ("机构本地模型（年龄三分位三机构）", "0.8362-0.9351"),
    ("联邦学习 FedAvg（数据不出院）", "0.9091"),
    ("集中训练上界（现实不可行）", "0.9019"),
]
tbl = s.shapes.add_table(4, 2, Inches(0.7), Inches(1.6), Inches(6.2), Inches(3.2)).table
tbl.columns[0].width = Inches(4.2)
tbl.columns[1].width = Inches(2.0)
for r, (a, b) in enumerate(rows):
    tbl.cell(r, 0).text = a
    tbl.cell(r, 1).text = b
    for c in range(2):
        cell = tbl.cell(r, c)
        cell.text_frame.paragraphs[0].runs[0].font.size = Pt(14)
        if r in (0, 2):
            cell.text_frame.paragraphs[0].runs[0].font.bold = True
bullets(s, 7.4, 1.8, 5.4, 4.3, [
    "真实患者冠脉病变终点：联邦追平集中上界，不依赖合成数据假设",
    "高于公开文献同数据集典型基线（0.85-0.90）",
    "合成三院基准：联邦 0.7018 同样追平集中上界 0.7012",
    "种子固定：GET /api/federation/benchmark 随时复现",
    "工程底座：445 个单元测试 + CI/CD + 三级降级容错",
], size=15)
text(s, 0.7, 6.3, 12.0, 0.6,
     "双基准佐证：真实公开数据复验（0.9091）+ 合成流行病学基准（0.7018，阳性率贴近心衰 30 天再入院文献），"
     "一台普通笔记本（无 GPU）全本地运行，决赛现场零网络依赖",
     size=13, bold=True, color=CYAN)
page_no(s, 7)

# ---------- 8 商业模式 ----------
s = add_slide()
header(s, "商业模式：平台佣金 + 订阅 + 服务 + 监管 SaaS")
models = [
    ("平台佣金 20%", "数据产品交易额分成（分成架构已内置上线）"),
    ("模型 API 订阅", "联邦风险模型按年授权，¥80,000/年起"),
    ("治理服务费", "医院病历治理与数据资产化，按数据集计价"),
    ("监管 SaaS 年费", "面向卫健/数据局的合规看板"),
]
for i, (t, body) in enumerate(models):
    x = 0.7 + (i % 2) * 6.2
    y = 1.9 + (i // 2) * 1.75
    rect(s, x, y, 5.9, 1.45, RGBColor(0xF1, 0xF5, 0xF9))
    rect(s, x, y, 0.12, 1.45, CYAN)
    text(s, x + 0.35, y + 0.15, 5.4, 0.5, t, size=18, bold=True, color=CYAN)
    text(s, x + 0.35, y + 0.75, 5.4, 0.6, body, size=14)
text(s, 0.7, 5.7, 12.0, 1.2,
     "差异化：不做「又一个医疗大模型」，做大模型时代的数据要素基础设施。\n"
     "与医疗 AI 应用共生（我们的应用生态就是数据买方）；与隐私计算厂商相比多出「流通交易 + 监管合规」完整闭环。",
     size=15, color=DARK)
page_no(s, 8)

# ---------- 9 落地计划 ----------
s = add_slide()
header(s, "温州落地计划：智能谷起步，医共体试点")
steps = [
    ("T+3月", "入驻中国（温州）智能谷；注册项目公司；申请鹿城区「AI+医疗应用场景示范项目库」"),
    ("T+6月", "与 1 家三甲 + 2 家基层机构签约，真实医共体环境联邦建模试点"),
    ("T+12月", "数据产品目录上线运营，首笔真实数据要素交易；申报省标杆项目（衔接最高 500 万奖补）"),
    ("T+24月", "复制到温州其他区县及浙南医共体，建成区域医疗数据要素流通基础设施"),
]
for i, (t, body) in enumerate(steps):
    y = 1.8 + i * 1.25
    rect(s, 0.7, y, 1.5, 1.0, CYAN if i % 2 == 0 else DARK)
    text(s, 0.8, y + 0.25, 1.3, 0.6, t, size=18, bold=True, color=WHITE)
    rect(s, 2.2, y, 10.4, 1.0, RGBColor(0xF8, 0xFA, 0xFC))
    text(s, 2.5, y + 0.2, 9.9, 0.7, body, size=15)
text(s, 0.7, 6.9, 12.0, 0.5,
     "政策衔接：重大科技攻关最高 200 万 · 软件开发最高 200 万 · 算力券/模型券/语料券每年最高 50 万 · 领军型创业项目 30-3000 万",
     size=13, color=GRAY)
page_no(s, 9)

# ---------- 10 团队与知产 ----------
s = add_slide()
header(s, "团队与知识产权")
bullets(s, 0.8, 1.9, 6.0, 4.5, [
    "三人核心团队：全栈开发 / 审核文档 / 数据库工程（详见报名表）",
    "AI 辅助开发范式：工程效率即竞争壁垒，50,809 行代码三周完成",
    "工程底座成熟：445 个单元测试、CI/CD、Docker 一键部署、魔搭云端部署经验",
], size=16)
bullets(s, 7.0, 1.9, 5.8, 4.5, [
    "软件著作权 ×3（申报中）：\n  平台 V1.0 / 治理 Copilot / 联邦引擎",
    "发明专利交底书 ×1：\n  联邦统计与差分隐私分级协同的医疗数据流通方法",
    "合规对齐：数据安全法 / 个保法 / 数据二十条",
    "演示数据全部为合成数据，零合规风险",
], size=16)
page_no(s, 10)

# ---------- 11 结尾 ----------
s = add_slide()
rect(s, 0, 0, 13.333, 7.5, DARK)
text(s, 1.0, 2.4, 11.3, 1.0, "让沉睡的医疗数据，", size=36, bold=True, color=WHITE)
text(s, 1.0, 3.3, 11.3, 1.0, "变成可确权、可定价、可流通、可监管的数据要素。", size=36, bold=True, color=CYAN)
text(s, 1.0, 5.2, 11.3, 0.6, "瓯医数链 OuMedTrust · 医疗数据 · 可用不可见", size=18, color=RGBColor(0x94, 0xA3, 0xB8))
text(s, 1.0, 6.3, 11.3, 0.6, "第二届全球技术创新大赛 AI+医疗专题赛 · 2026 中国温州", size=13, color=RGBColor(0x94, 0xA3, 0xB8))

prs.save(OUT)
print(f"PPT 已生成: {OUT}（{len(prs.slides.__iter__.__self__._sldIdLst)} 页）")
