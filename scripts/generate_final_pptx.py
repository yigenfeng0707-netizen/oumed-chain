# -*- coding: utf-8 -*-
"""瓯医数链决赛路演 PPT（15 页故事版 · 视觉版同款视觉系统）

视觉系统与《初赛路演-视觉版》(ppt_assemble.py) 完全一致：
  全屏场景图 + 青色大标题 + 描边圆角卡片 + 描边大数字卡 + 底部「▎」结论条
内容为 2026-09-02 决赛事实口径（0.9091 双基准 / 支付宝 live 闭环 / 448 测试 / 软著双轨）。

故事主线：一间病房的处方数据，为什么救不了隔壁医院的病人？
→ 三重困境 → 四层架构 → 数据旅程 → 双基准实证 → 真实运行（支付闭环）→ 商业与落地 → 愿景
运行：backend/.venv/Scripts/python scripts/generate_final_pptx.py
输出：docs/瓯医数链-决赛路演.pptx
"""
import os

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
from pptx.util import Inches, Pt

ROOT = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
CH = os.path.join(ROOT, "docs", "PPT素材", "charts")
SC = os.path.join(ROOT, "docs", "PPT素材", "scenes")
SHOT = os.path.join(ROOT, "docs", "screenshots")
OUT = os.path.join(ROOT, "docs", "瓯医数链-决赛路演.pptx")

# ---- 设计令牌（与视觉版一致）----
BG = RGBColor(0x0A, 0x1A, 0x3F)
CYAN = RGBColor(0x00, 0xD4, 0xFF)
ORANGE = RGBColor(0xFF, 0x6B, 0x35)
GOLD = RGBColor(0xFF, 0xD7, 0x00)
GREEN = RGBColor(0x00, 0xC8, 0x53)
BLUE = RGBColor(0x1E, 0x5B, 0xFF)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GREY = RGBColor(0xA9, 0xBA, 0xDF)
CARD = RGBColor(0x12, 0x23, 0x4E)
BAND = RGBColor(0x0E, 0x1F, 0x4A)
FONT = "微软雅黑"

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]
PAGE = [0]


def _font(r, size, color, bold):
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = color
    r.font.name = FONT
    rPr = r._r.get_or_add_rPr()
    for tag in ("a:ea", "a:cs"):
        el = rPr.find(qn(tag))
        if el is None:
            el = rPr.makeelement(qn(tag), {})
            rPr.append(el)
        el.set("typeface", FONT)


def text(s, x, y, w, h, content, size=16, color=WHITE, bold=False,
         align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, line_spacing=1.15):
    tb = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    for i, line in enumerate(content.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = line_spacing
        r = p.add_run()
        r.text = line
        _font(r, size, color, bold)
    return tb


def bullets(s, x, y, w, h, items, size=14.5, gap=10, color=GREY):
    tb = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    for i, it in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(gap)
        p.line_spacing = 1.2
        r = p.add_run()
        r.text = "▪ " + it
        _font(r, size, color, False)
    return tb


def add_card(s, x, y, w, h, line_color, fill=CARD):
    box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                             Inches(x), Inches(y), Inches(w), Inches(h))
    box.fill.solid()
    box.fill.fore_color.rgb = fill
    box.line.color.rgb = line_color
    box.line.width = Pt(1.5)
    box.shadow.inherit = False
    return box


def semi_transparent(shape, pct):
    """给纯色填充加 alpha（pct=不透明度百分比*1000）"""
    sp = shape.fill._xPr.find(qn("a:solidFill"))
    clr = sp.find(qn("a:srgbClr"))
    alpha = clr.makeelement(qn("a:alpha"), {"val": str(pct)})
    clr.append(alpha)


def metric_box(s, x, y, w, h, value, label, color, value_size=30,
               label_size=13):
    add_card(s, x, y, w, h, color)
    text(s, x + 0.1, y + 0.16, w - 0.2, h * 0.55, value, value_size, color,
         bold=True, align=PP_ALIGN.CENTER)
    text(s, x + 0.1, y + h - 0.62, w - 0.2, 0.55, label, label_size, GREY,
         align=PP_ALIGN.CENTER)


def page_header(s, title, subtitle=None):
    PAGE[0] += 1
    sl = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width,
                            Inches(0.09))
    sl.fill.solid()
    sl.fill.fore_color.rgb = CYAN
    sl.line.fill.background()
    sl.shadow.inherit = False
    text(s, 0.55, 0.35, 11.5, 0.9, title, 29, CYAN, bold=True)
    if subtitle:
        text(s, 0.55, 1.18, 11.5, 0.5, subtitle, 14.5, GREY)
    text(s, 12.2, 0.42, 0.8, 0.4, f"{PAGE[0]:02d}", 12, GREY,
         align=PP_ALIGN.RIGHT)


def takeaway(s, t, color=CYAN):
    add_card(s, 0.55, 6.68, 12.2, 0.62, color, fill=BAND)
    text(s, 0.85, 6.68, 11.7, 0.62, "▎" + t, 15, WHITE, bold=True,
         anchor=MSO_ANCHOR.MIDDLE)


def full_picture(s, path):
    s.shapes.add_picture(path, 0, 0, width=prs.slide_width,
                         height=prs.slide_height)


def pic(s, path, x, y, h=None, w=None):
    if os.path.exists(path):
        if h is not None:
            s.shapes.add_picture(path, Inches(x), Inches(y), height=Inches(h))
        else:
            s.shapes.add_picture(path, Inches(x), Inches(y), width=Inches(w))
        return True
    return False


def chart(name):
    return os.path.join(CH, name)


def shot(name):
    return os.path.join(SHOT, name)


def scene(name):
    return os.path.join(SC, name)


# ---------- 1 封面：全屏 hero + 底部横幅（视觉版同款） ----------
s = prs.slides.add_slide(BLANK)
full_picture(s, scene("cover_hero.jpg"))
band = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(4.9),
                          prs.slide_width, Inches(2.6))
band.fill.solid()
band.fill.fore_color.rgb = BG
band.line.fill.background()
band.shadow.inherit = False
text(s, 0.8, 5.05, 12, 0.95, "瓯医数链 OuMedTrust", 44, WHITE, bold=True)
text(s, 0.8, 6.05, 12, 0.55,
     "医疗数据要素可信流通平台 · 让数据「可用不可见、可控可计量」", 20, CYAN)
text(s, 0.8, 6.72, 12, 0.5,
     "第二届全球技术创新大赛 · AI+医疗专题赛 · 赛道二：医疗大模型与数据", 13, GREY)
text(s, 10.7, 0.4, 2.3, 0.5, "数据不出院 · 价值出院门", 13, GOLD, bold=True,
     align=PP_ALIGN.RIGHT)

# ---------- 2 故事开场：全屏孤岛场景 ----------
s = prs.slides.add_slide(BLANK)
full_picture(s, scene("pain_silos.jpg"))
ov = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width,
                        prs.slide_height)
ov.fill.solid()
ov.fill.fore_color.rgb = BG
ov.line.fill.background()
ov.shadow.inherit = False
semi_transparent(ov, 62000)
text(s, 0.55, 0.5, 12, 0.9, "一间病房的处方数据，为什么救不了隔壁医院的病人？",
     29, CYAN, bold=True)
text(s, 0.7, 1.8, 7.6, 2.6,
     "同一座城市，三家医院。\n"
     "三甲医院见过 4200 例心衰患者，县医院 2400 例，社区卫生中心 1100 例。\n"
     "隔壁医院的治疗经验，救不了这一位病人——数据被锁在各自的数据库里。",
     18, WHITE, line_spacing=1.5)
text(s, 0.7, 4.4, 7.6, 1.6,
     "不是医院不想共享，而是：法律不敢、技术不会、市场没有。\n"
     "「数据大池化」在《数据安全法》《个人信息保护法》面前已经走不通。",
     15.5, GREY, line_spacing=1.4)
takeaway(s, "瓯医数链的答案：数据不出院，价值出院门。", CYAN)

# ---------- 3 三重困境：描边圆角卡 ----------
s = prs.slides.add_slide(BLANK)
page_header(s, "三重困境：医疗数据动不了、不值钱、没基建")
cols = [
    ("不敢共享", CYAN,
     "数据敏感、权属模糊\n数据安全法/个保法硬约束\n传统大池化模式走不通"),
    ("不会增值", ORANGE,
     "基层数据缺失率 12%-32%\n单院建模 AUC 仅 0.69\n数据资产沉睡、无法变现"),
    ("没有基建", GOLD,
     "定价、授权、分成、监管\n缺乏「可用不可见」的\n技术底座与合规流程"),
]
for i, (t, c, body) in enumerate(cols):
    x = 0.7 + i * 4.15
    add_card(s, x, 1.75, 3.85, 2.6, c)
    text(s, x + 0.28, 1.98, 3.3, 0.6, t, 22, c, bold=True)
    text(s, x + 0.28, 2.72, 3.3, 1.5, body, 14.5, GREY, line_spacing=1.35)
add_card(s, 0.7, 4.75, 12.2, 1.5, CYAN, fill=BAND)
text(s, 1.05, 4.95, 11.6, 1.15,
     "政策风口：数据要素×医疗健康为国家数据局重点行动方向；浙江省是要素市场化改革试点；\n"
     "本赛道主办方第一名 = 温州市经信局，数据要素正是其核心 KPI。",
     15, WHITE, line_spacing=1.35)
takeaway(s, "数据不敢动 → AI 吃不到跨院数据 → 患者享受不到更好的模型", ORANGE)

# ---------- 4 四层架构 ----------
s = prs.slides.add_slide(BLANK)
page_header(s, "解决方案：治理 - 协作 - 流通 - 监管，四层完整闭环",
            "闭环故事：治理脱敏 → 上架市场 → 联邦消费 → 交易存证 → 收益反哺")
pic(s, chart("05_architecture.png"), 1.35, 1.8, h=4.75)
takeaway(s, "不做「又一个医疗大模型」，做数据要素基础设施 —— 与医疗 AI 应用共生")

# ---------- 5 数据旅程：五步描边卡 ----------
s = prs.slides.add_slide(BLANK)
page_header(s, "一条患者数据的旅程：从社区病历到要素市场")
steps = [
    ("1. 治理", "社区病历 AI 脱敏\n结构化上架", CYAN),
    ("2. 协作", "联邦消费三院数据\n数据不出院建模", BLUE),
    ("3. 流通", "模型被 Agent\n授权采购", ORANGE),
    ("4. 结算", "70/20/10\n自动分成", GOLD),
    ("5. 监管", "sha256 存证\n看板实时可见", GREEN),
]
for i, (t, body, c) in enumerate(steps):
    x = 0.55 + i * 2.52
    add_card(s, x, 2.3, 2.25, 2.6, c)
    text(s, x + 0.2, 2.52, 1.9, 0.55, t, 18, c, bold=True)
    text(s, x + 0.2, 3.25, 1.9, 1.5, body, 13.5, GREY, line_spacing=1.35)
    if i < 4:
        text(s, x + 2.2, 3.35, 0.4, 0.6, "→", 24, CYAN, bold=True)
takeaway(s, "这不是流程图——每一步都已经在线上真实跑通（后 5 页逐一验证）")

# ---------- 6 联邦学习实测 ----------
s = prs.slides.add_slide(BLANK)
page_header(s, "联邦学习实测：94 毫秒，数据不出院完成三院联合建模",
            "三家异构医院（三甲 4200 / 县 2400 / 社区 1100 例，特征缺失 5%-22%）")
pic(s, shot("02-联邦协作网络-训练结果.png"), 0.55, 1.85, w=7.55)
metric_box(s, 8.4, 1.95, 4.35, 1.35, "0.7018", "合成三院联邦 AUC\n超过任何单院模型",
           CYAN)
metric_box(s, 8.4, 3.5, 4.35, 1.35, "0.7012", "集中训练上界\n追平（现实不可行）",
           ORANGE)
metric_box(s, 8.4, 5.05, 4.35, 1.35, "94ms", "单轮联邦聚合\n纯 CPU 笔记本可复现",
           GREEN)
takeaway(s, "逐院公平性：三家全部获益，基层提升最大 —— 差分隐私仅损失 0.01 AUC")

# ---------- 7 真实数据 0.9091（核心页） ----------
s = prs.slides.add_slide(BLANK)
page_header(s, "真实患者数据：AUC 0.9091 追平大池化训练上界",
            "UCI 心脏病真实队列（Cleveland）· 297 例真实患者 · 非 IID 三机构（按年龄三分位）")
pic(s, chart("01_auc_real.png"), 0.55, 1.85, w=9.3)
metric_box(s, 10.15, 2.1, 2.7, 1.35, "0.9091", "联邦 AUC\n数据不出院", CYAN)
metric_box(s, 10.15, 3.65, 2.7, 1.35, "0.9019", "集中上界\n(现实不可行)", ORANGE)
metric_box(s, 10.15, 5.2, 2.7, 1.35, "297", "真实患者\nCleveland 队列", GREEN)
takeaway(s, "高于公开文献同数据集典型基线(0.85-0.90) —— 联邦增益在真实分布下成立",
         GREEN)

# ---------- 8 AI 病历治理 ----------
s = prs.slides.add_slide(BLANK)
page_header(s, "AI 病历治理 Copilot：院内网脱敏结构化，永不宕机",
            "本地 qwen3:4b 推理 · 4 类 PHI 零漏检 · LLM 失效自动降级规则引擎")
pic(s, shot("04-AI病历治理-脱敏对比.png"), 0.55, 1.85, w=7.55)
metric_box(s, 8.4, 1.95, 4.35, 1.35, "0 漏检", "身份证/手机号/姓名/住院号\n4 类 PHI 全检出",
           ORANGE)
metric_box(s, 8.4, 3.5, 4.35, 1.35, "7 字段", "非结构化病历 → 标准 JSON\n单份治理 10-60 秒",
           CYAN)
metric_box(s, 8.4, 5.05, 4.35, 1.35, "断网可跑", "云端密钥零依赖\n治理产物直通要素市场",
           GREEN)
takeaway(s, "治理是流通的前提：先脱敏结构化，才谈得上确权与定价")

# ---------- 9 支付宝真实收款闭环 ----------
s = prs.slides.add_slide(BLANK)
page_header(s, "已实现支付宝真实交易闭环：不是沙箱，是真金白银",
            "2026-09-02 终验 · 订单 OM260902105109807F · 下单到落账 4 分 42 秒")
pic(s, shot("alipay_live_cashier_checkout.png"), 0.55, 1.85, w=7.55)
metric_box(s, 8.4, 1.95, 4.35, 1.35, "¥3.90", "真实订单第三方扫码实付\npending → paid",
           CYAN)
metric_box(s, 8.4, 3.5, 4.35, 1.35, "4分42秒", "下单 → 落账 → 存证上链\n全程实测",
           ORANGE)
metric_box(s, 8.4, 5.05, 4.35, 1.35, "20260902…789", "真实交易号回写\n管理端对账实时可见",
           GREEN, value_size=24)
takeaway(s, "支付宝开放平台应用过审上线，生产收银台直达 —— 评委可现场登录管理后台核对",
         GOLD)

# ---------- 10 监管看板 ----------
s = prs.slides.add_slide(BLANK)
page_header(s, "监管方实时可见：流通金额 · 收益分配 · 隐私事件 0 起",
            "审计存证链 sha256 串联，任何篡改导致链条断裂，一键校验")
pic(s, shot("06-数据要素市场-监管看板.png"), 0.55, 1.85, w=7.55)
metric_box(s, 8.4, 1.95, 4.35, 1.35, "70/20/10", "医院 / 平台 / 贡献者\n自动分成",
           CYAN)
metric_box(s, 8.4, 3.5, 4.35, 1.35, "0 起", "隐私事件\n全程留痕 · 实时监控", GREEN)
metric_box(s, 8.4, 5.05, 4.35, 1.35, "4 法规", "数据安全法/个保法\n数据二十条/GB&T 39725",
           GOLD, value_size=26)
takeaway(s, "每一笔交易自动存证上链 —— 篡改即断链，监管方一键校验")

# ---------- 11 工程实证：六宫格大数字 ----------
s = prs.slides.add_slide(BLANK)
page_header(s, "工程实证：一个已经在公网真实运行、能收款、能扛并发的系统")
grid = [
    ("448 项", "单元测试全部通过\n含生产鉴权与越权审计回归", CYAN),
    ("100%", "LLM 全链路压测成功率\n5 路 15/15 · 10 路 20/20", GREEN),
    ("4分42秒", "支付宝真实收款闭环\n下单 → 落账 → 存证上链", ORANGE),
    ("2 模型", "主备双模高可用\n429 限频自动切换实战验证", BLUE),
    ("0 起", "隐私事件\n全程留痕 · 存证链校验", GOLD),
    ("24/7", "魔搭线上 Demo 全时可访问\nms.show 公网直达", CYAN),
]
for i, (v, l, c) in enumerate(grid):
    x = 0.7 + (i % 3) * 4.15
    y = 1.75 + (i // 3) * 2.15
    metric_box(s, x, y, 3.85, 1.9, v, l, c, value_size=32)
text(s, 0.7, 6.1, 12.0, 0.5,
     "线上 Demo：https://gsym236998-oumed-chain-demo.ms.show（9 智能体 · 泛癌卫士 · 管理后台）",
     14, GREY)
takeaway(s, "系统已过终验：能收款、能扛并发、能查证 —— 不是 PPT 工程")

# ---------- 12 商业模式 ----------
s = prs.slides.add_slide(BLANK)
page_header(s, "商业模式：四大收入源，应用生态就是数据买方",
            "区域医共体(三甲+县医院+社区) · 卫健/数据局 · 保险精算 · 药研 CRO")
pic(s, chart("04_revenue.png"), 0.7, 1.85, w=6.4)
diffs = [
    "模型 API 订阅：¥80,000/年起（目录已上架）",
    "平台佣金：交易额 20%（分成已内置）",
    "治理服务费：按数据集计价（Copilot 已上线）",
    "监管 SaaS：卫健/数据局看板年费",
]
for i, d in enumerate(diffs):
    add_card(s, 7.4, 1.95 + i * 1.06, 5.4, 0.9, CYAN if i == 0 else CARD)
    text(s, 7.68, 1.95 + i * 1.06, 4.9, 0.9, d, 14,
         WHITE if i == 0 else GREY, bold=(i == 0),
         anchor=MSO_ANCHOR.MIDDLE)
takeaway(s, "不做「又一个医疗大模型」，做大模型时代的数据要素基础设施")

# ---------- 13 温州落地 ----------
s = prs.slides.add_slide(BLANK)
page_header(s, "温州落地：鹿城起笔，浙南成网",
            "衔接温州市经信局核心 KPI · 中国（温州）智能谷天然载体")
pic(s, chart("07_timeline.png"), 1.0, 1.8, w=11.3)
takeaway(s, "落地后可申报：重大科技攻关 200 万 · 场景应用奖补 500 万 · 领军创业 30-3000 万",
         GOLD)

# ---------- 14 团队与知产 ----------
s = prs.slides.add_slide(BLANK)
page_header(s, "团队与知识产权：权属清晰，双轨并进")
add_card(s, 0.7, 1.8, 5.9, 4.4, CYAN)
text(s, 1.0, 2.05, 5.3, 0.55, "团队", 20, CYAN, bold=True)
bullets(s, 1.0, 2.75, 5.3, 3.3, [
    "三人核心团队：全栈开发 / 审核文档 / 数据库工程",
    "AI 辅助开发范式：工程效率即竞争壁垒",
    "448 项测试 + CI/CD + Docker 一键部署",
    "PIA 个人信息保护影响评估报告已产出",
], size=14.5, color=GREY)
add_card(s, 6.85, 1.8, 5.9, 4.4, GOLD)
text(s, 7.15, 2.05, 5.3, 0.55, "知识产权", 20, GOLD, bold=True)
bullets(s, 7.15, 2.75, 5.3, 3.3, [
    "软著主轨：瓯医数链 ×3 材料齐备提交中",
    "已登记软著 ×2（个人独有，证书原件核验）：\n"
    "   视频人数统计（视觉 AI）· 充电桩故障检测（时序异常）",
    "共有软著 ×1（技术能力佐证）",
    "发明专利在案：《数据交易的管理方法…》\n本人为发明人（申请号 2025111566552）",
], size=14, color=GREY)
takeaway(s, "自主知识产权权属链条清晰 —— 决赛资格与商业化授权均无瑕疵")

# ---------- 15 结尾：全屏愿景 + 半透明横幅 ----------
s = prs.slides.add_slide(BLANK)
full_picture(s, scene("vision_wenzhou.jpg"))
band = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(4.6),
                          prs.slide_width, Inches(2.9))
band.fill.solid()
band.fill.fore_color.rgb = BG
band.line.fill.background()
band.shadow.inherit = False
semi_transparent(band, 80000)
text(s, 0.8, 4.72, 12, 0.8, "让沉睡的医疗数据，", 32, WHITE, bold=True)
text(s, 0.8, 5.42, 12, 0.7, "变成可确权、可定价、可流通、可监管的数据要素。",
     32, CYAN, bold=True)
text(s, 0.8, 6.28, 12, 0.5,
     "448 项单元测试全绿 · 支付宝真实收款闭环 · 软著双轨并进 · 线上 Demo 全时可访问",
     15, CYAN)
text(s, 0.8, 6.82, 12, 0.5, "瓯医数链 OuMedTrust · 恳请评委指正 · 2026 中国温州",
     13, GREY)

prs.save(OUT)
print("PPT saved:", OUT, "slides:", len(prs.slides._sldIdLst))
