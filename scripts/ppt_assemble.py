# -*- coding: utf-8 -*-
"""瓯医数链初赛路演 PPT 组装（12 页叙事版，2026-08-31 真实数据口径）"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

BASE = os.path.join(os.path.dirname(__file__), "..", "docs", "PPT素材")
CH = os.path.join(BASE, "charts")
SC = os.path.join(BASE, "scenes")
OUT_PPTX = os.path.join(os.path.dirname(__file__), "..", "docs",
                        "瓯医数链-初赛路演-视觉版.pptx")

BG = RGBColor(0x0A, 0x1A, 0x3F)
CYAN = RGBColor(0x00, 0xD4, 0xFF)
ORANGE = RGBColor(0xFF, 0x6B, 0x35)
GOLD = RGBColor(0xFF, 0xD7, 0x00)
GREEN = RGBColor(0x00, 0xC8, 0x53)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GREY = RGBColor(0xA9, 0xBA, 0xDF)
CARD = RGBColor(0x12, 0x23, 0x4E)
FONT = "微软雅黑"

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]


def set_bg(slide, color=BG):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_text(slide, x, y, w, h, text, size, color=WHITE, bold=False,
             align=PP_ALIGN.LEFT, font=FONT, anchor=MSO_ANCHOR.TOP,
             line_spacing=1.15):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    for i, line in enumerate(text.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = line_spacing
        r = p.add_run()
        r.text = line
        r.font.size = Pt(size)
        r.font.color.rgb = color
        r.font.bold = bold
        r.font.name = font
    return tb


def add_card(slide, x, y, w, h, line_color, fill=CARD):
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                 Inches(x), Inches(y), Inches(w), Inches(h))
    box.fill.solid()
    box.fill.fore_color.rgb = fill
    box.line.color.rgb = line_color
    box.line.width = Pt(1.5)
    box.shadow.inherit = False
    return box


def semi_transparent(shape, pct):
    """给纯色填充加 alpha（pct=不透明度百分比*1000）"""
    sp = shape.fill._xPr.find(qn('a:solidFill'))
    clr = sp.find(qn('a:srgbClr'))
    alpha = clr.makeelement(qn('a:alpha'), {'val': str(pct)})
    clr.append(alpha)


def metric_box(slide, x, y, w, h, value, label, color, value_size=32,
               label_size=13):
    add_card(slide, x, y, w, h, color)
    add_text(slide, x + 0.1, y + 0.18, w - 0.2, h * 0.55, value,
             value_size, color, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, x + 0.1, y + h - 0.62, w - 0.2, 0.55, label,
             label_size, GREY, align=PP_ALIGN.CENTER)


def page_header(slide, title, subtitle=None, page=None):
    add_text(slide, 0.55, 0.35, 11.5, 0.9, title, 30, CYAN, bold=True)
    if subtitle:
        add_text(slide, 0.55, 1.15, 11.5, 0.5, subtitle, 15, GREY)
    if page:
        add_text(slide, 12.2, 0.42, 0.8, 0.4, page, 12, GREY,
                 align=PP_ALIGN.RIGHT)


def takeaway(slide, text, color=CYAN):
    add_card(slide, 0.55, 6.7, 12.2, 0.62, color,
             fill=RGBColor(0x0E, 0x1F, 0x4A))
    add_text(slide, 0.85, 6.78, 11.7, 0.5, "▎" + text, 15, WHITE, bold=True,
             anchor=MSO_ANCHOR.MIDDLE)


def full_picture(slide, path):
    slide.shapes.add_picture(path, 0, 0, width=prs.slide_width,
                             height=prs.slide_height)


# ---------- S1 封面 ----------
s = prs.slides.add_slide(BLANK)
set_bg(s)
full_picture(s, os.path.join(SC, "cover_hero.jpg"))
band = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(4.9),
                          prs.slide_width, Inches(2.6))
band.fill.solid()
band.fill.fore_color.rgb = BG
band.line.fill.background()
band.shadow.inherit = False
add_text(s, 0.8, 5.05, 12, 0.95, "瓯医数链 OuMedTrust", 44, WHITE, bold=True)
add_text(s, 0.8, 6.05, 12, 0.55, "医疗数据要素可信流通平台 · 让数据「可用不可见、可控可计量」",
         20, CYAN)
add_text(s, 0.8, 6.72, 12, 0.5,
         "第二届全球技术创新大赛 · AI+医疗专题赛 · 赛道二：医疗大模型与数据",
         13, GREY)
add_text(s, 10.9, 0.4, 2.0, 0.5, "数据不出院 · 价值出院门", 13, GOLD,
         bold=True, align=PP_ALIGN.RIGHT)

# ---------- S2 痛点 ----------
s = prs.slides.add_slide(BLANK)
set_bg(s)
full_picture(s, os.path.join(SC, "pain_silos.jpg"))
ov = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width,
                        prs.slide_height)
ov.fill.solid()
ov.fill.fore_color.rgb = BG
ov.line.fill.background()
ov.shadow.inherit = False
semi_transparent(ov, 55000)
add_text(s, 0.55, 0.5, 11, 0.9, "三重困局：医院数据沉睡在孤岛上", 30, CYAN,
         bold=True)
cards = [
    ("不敢共享", "数据安全法/个保法硬约束\n权属模糊，数据大池化模式出局", ORANGE),
    ("不会增值", "基层特征缺失 12%-32%\n单院建模 AUC 仅 0.69，资产无法变现", ORANGE),
    ("没有基建", "定价、授权、分成、监管\n缺乏「可用不可见」底座与合规流程", ORANGE),
]
for i, (t, d, c) in enumerate(cards):
    x = 0.7 + i * 4.15
    add_card(s, x, 4.6, 3.85, 2.2, c)
    add_text(s, x + 0.25, 4.78, 3.4, 0.6, t, 22, c, bold=True)
    add_text(s, x + 0.25, 5.42, 3.4, 1.25, d, 13.5, GREY)
takeaway(s, "数据不敢动 → AI 吃不到跨院数据 → 患者享受不到更好的模型", ORANGE)

# ---------- S3 使命定位 ----------
s = prs.slides.add_slide(BLANK)
set_bg(s)
page_header(s, "我们的答案：数据要素的「交易所 + 治理工厂」", None, "02")
add_text(s, 0.55, 1.6, 12.2, 1.2,
         "让沉睡在医院服务器里的数据，\n变成可确权、可定价、可流通、可监管的数据要素。",
         26, WHITE, bold=True, line_spacing=1.3)
boxes = [
    ("可用不可见", "联邦学习 · 数据不出院\n跨院联合建模", CYAN),
    ("可控可计量", "用途限定授权 · 差分隐私\n全程审计存证", GOLD),
    ("可控流通增值", "数据产品目录 · 收益分成\n监管实时看板", GREEN),
]
for i, (t, d, c) in enumerate(boxes):
    x = 0.7 + i * 4.15
    add_card(s, x, 3.6, 3.85, 2.35, c)
    add_text(s, x + 0.25, 3.82, 3.4, 0.6, t, 21, c, bold=True)
    add_text(s, x + 0.25, 4.5, 3.4, 1.3, d, 14, GREY)
takeaway(s, "1 个可信数据底座 + N 个医疗 AI 应用生态（9 个智能体即数据消费方）")

# ---------- S4 架构 ----------
s = prs.slides.add_slide(BLANK)
set_bg(s)
page_header(s, "四层架构闭环：治理 → 协作 → 流通 → 监管",
            "闭环故事：治理脱敏 → 上架市场 → 联邦消费 → 交易存证 → 收益反哺", "03")
s.shapes.add_picture(os.path.join(CH, "05_architecture.png"),
                     Inches(0.55), Inches(1.75), width=Inches(12.2))
takeaway(s, "不做「又一个医疗大模型」，做数据要素基础设施 —— 与医疗 AI 应用共生")

# ---------- S5 真实数据验证（核心页） ----------
s = prs.slides.add_slide(BLANK)
set_bg(s)
page_header(s, "真实患者数据：AUC 0.9091 追平大池化训练上界",
            "UCI 心脏病真实队列 · 297 例真实患者 · 非 IID 三机构（按年龄三分位）", "04")
s.shapes.add_picture(os.path.join(CH, "01_auc_real.png"),
                     Inches(0.55), Inches(1.8), width=Inches(9.3))
metric_box(s, 10.15, 2.1, 2.7, 1.35, "0.9091", "联邦 AUC\n数据不出院", CYAN,
           value_size=30)
metric_box(s, 10.15, 3.65, 2.7, 1.35, "0.9019", "集中上界\n(现实不可行)", ORANGE,
           value_size=30)
metric_box(s, 10.15, 5.2, 2.7, 1.35, "297", "真实患者\nCleveland 队列", GREEN,
           value_size=30)
takeaway(s, "高于公开文献同数据集典型基线(0.85-0.90) —— 联邦增益在真实分布下成立", GREEN)

# ---------- S6 收敛曲线 ----------
s = prs.slides.add_slide(BLANK)
set_bg(s)
page_header(s, "第 2 轮即收敛：94ms 一轮聚合，笔记本 CPU 可复现",
            "种子固定 · 12 轮稳定 · 含差分隐私噪声", "05")
s.shapes.add_picture(os.path.join(CH, "02_convergence.png"),
                     Inches(0.55), Inches(1.8), width=Inches(9.3))
metric_box(s, 10.15, 2.1, 2.7, 1.35, "94ms", "单轮联邦聚合\n实测耗时", CYAN,
           value_size=30)
metric_box(s, 10.15, 3.65, 2.7, 1.35, "12轮", "收敛轮次\n全程可复现", GOLD,
           value_size=30)
metric_box(s, 10.15, 5.2, 2.7, 1.35, "0起", "隐私事件\n全程监控", GREEN,
           value_size=30)
takeaway(s, "轻噪声差分隐私仅损失 0.01 AUC —— 隐私与效用可分档权衡", CYAN)

# ---------- S7 AI 治理 Copilot ----------
s = prs.slides.add_slide(BLANK)
set_bg(s)
page_header(s, "AI 病历治理 Copilot：非结构化病历 → 干净数据资产",
            "本地 qwen3:4b 院内网推理 · 断网可跑 · LLM 失效自动降级规则引擎", "06")
steps = [
    ("原始病历", "非结构化文本\n自由格式", RGBColor(0x5B, 0x7B, 0xC0)),
    ("PHI 脱敏", "身份证/手机号/姓名\n住院号 4 类零漏检", ORANGE),
    ("结构化", "7 字段 JSON\n患者/主诉/诊断/用药…", CYAN),
    ("数据资产", "零 PHI 残留\n可上架数据要素市场", GREEN),
]
for i, (t, d, c) in enumerate(steps):
    x = 0.55 + i * 3.25
    add_card(s, x, 1.85, 2.85, 1.95, c)
    add_text(s, x + 0.22, 2.0, 2.45, 0.55, t, 18, c, bold=True)
    add_text(s, x + 0.22, 2.62, 2.45, 1.1, d, 12.5, GREY)
    if i < 3:
        ar = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW,
                                Inches(x + 2.9), Inches(2.65),
                                Inches(0.3), Inches(0.35))
        ar.fill.solid()
        ar.fill.fore_color.rgb = CYAN
        ar.line.fill.background()
        ar.shadow.inherit = False
metric_box(s, 0.55, 4.35, 3.0, 1.5, "0", "PHI 漏检\n4 类实体全检出", ORANGE)
metric_box(s, 3.9, 4.35, 3.0, 1.5, "7", "结构化字段\n一键 JSON", CYAN)
metric_box(s, 7.25, 4.35, 3.0, 1.5, "10-60s", "单份病历治理\n本地推理实测", GOLD)
metric_box(s, 10.6, 4.35, 2.25, 1.5, "永不", "宕机\nLLM 失效自动降级", GREEN)
takeaway(s, "治理是流通的前提：先脱敏结构化，才谈得上确权与定价")

# ---------- S8 三引擎真实数据 ----------
s = prs.slides.add_slide(BLANK)
set_bg(s)
page_header(s, "三引擎全部跑过真实公开数据集 —— 指标不是编的",
            "合成数据验证机制 + 真实公开数据复验双重佐证", "07")
engines = [
    ("联邦引擎", "UCI 心脏病\n297 例真实患者", "AUC 0.9091\n追平集中训练上界", CYAN),
    ("EEG 引擎", "PhysioNet eegmmidb\n8 条真实记录 · 7 名受试", "闭眼→放松 睁眼→疲劳\n符合脑电生理学常识", GREEN),
    ("影像引擎", "Oncoformer 泛癌预测\nCell 2026 · 温附医一院", "医院数据不出院\n模型作为产品流通", ORANGE),
]
for i, (t, src, res, c) in enumerate(engines):
    x = 0.55 + i * 4.25
    add_card(s, x, 1.9, 3.95, 3.9, c)
    badge = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                               Inches(x + 2.45), Inches(2.05),
                               Inches(1.3), Inches(0.4))
    badge.fill.solid()
    badge.fill.fore_color.rgb = c
    badge.line.fill.background()
    badge.shadow.inherit = False
    tf = badge.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = "REAL DATA"
    r.font.size = Pt(12)
    r.font.bold = True
    r.font.color.rgb = BG
    r.font.name = FONT
    add_text(s, x + 0.25, 2.1, 2.2, 0.5, t, 20, c, bold=True)
    add_text(s, x + 0.25, 2.9, 3.45, 1.1, src, 14, WHITE, line_spacing=1.3)
    add_text(s, x + 0.25, 4.25, 3.45, 1.4, res, 13.5, GREY, line_spacing=1.3)
takeaway(s, "五频段功率谱 · 医保政策联动 · ODC-By 许可署名 —— EEG 面板线上可点验", GREEN)

# ---------- S9 市场闭环 ----------
s = prs.slides.add_slide(BLANK)
set_bg(s)
page_header(s, "数据要素市场：每一笔交易可存证、可分成、可监管",
            "产品目录 → 用途限定授权 → 交易结算 → 收益自动分成", "08")
s.shapes.add_picture(os.path.join(CH, "03_sharing.png"),
                     Inches(0.55), Inches(1.8), width=Inches(4.6))
metric_box(s, 5.5, 2.3, 2.6, 1.35, "70/20/10", "医院/平台/贡献者\n自动分成", CYAN,
           value_size=24)
metric_box(s, 8.3, 2.3, 2.6, 1.35, "¥8万+", "联邦模型 API\n年订阅定价", GOLD,
           value_size=24)
metric_box(s, 11.0, 2.3, 1.85, 1.35, "0起", "隐私事件\n实时监控", GREEN)
s.shapes.add_picture(os.path.join(CH, "06_audit_chain.png"),
                     Inches(0.55), Inches(4.35), width=Inches(12.2))
takeaway(s, "审计存证链 sha256 串联上链 —— 篡改即断链，监管方一键校验")

# ---------- S10 商业模式 ----------
s = prs.slides.add_slide(BLANK)
set_bg(s)
page_header(s, "商业模式：四大收入源，服务医共体与监管两侧",
            "区域医共体(三甲+县医院+社区) · 卫健/数据局 · 保险精算 · 药研 CRO", "09")
s.shapes.add_picture(os.path.join(CH, "04_revenue.png"),
                     Inches(0.7), Inches(1.75), width=Inches(6.4))
diffs = [
    "vs 医疗大模型公司：我们做基础设施，应用生态即数据买方 —— 共生而非竞争",
    "vs 隐私计算厂商：多出「流通交易 + 监管合规」完整闭环",
    "vs 数据交易所：多出医院侧治理工厂与联邦技术底座",
]
for i, d in enumerate(diffs):
    add_card(s, 7.4, 1.9 + i * 1.45, 5.4, 1.25, CYAN if i == 0 else CARD)
    add_text(s, 7.65, 2.02 + i * 1.45, 5.0, 1.0, d, 13.5,
             WHITE if i == 0 else GREY, bold=(i == 0), line_spacing=1.25)
takeaway(s, "数据要素×医疗健康 = 国家数据局重点方向，浙江是市场化配置改革试点")

# ---------- S11 温州落地 ----------
s = prs.slides.add_slide(BLANK)
set_bg(s)
page_header(s, "温州落地：鹿城起笔，浙南成网",
            "衔接温州市经信局核心 KPI · 中国（温州）智能谷天然载体", "10")
s.shapes.add_picture(os.path.join(CH, "07_timeline.png"),
                     Inches(0.55), Inches(1.8), width=Inches(12.2))
takeaway(s, "落地后可申报：重大科技攻关 200 万 · 场景应用奖补 500 万 · 领军创业 30-3000 万", GOLD)

# ---------- S12 愿景致谢 ----------
s = prs.slides.add_slide(BLANK)
set_bg(s)
full_picture(s, os.path.join(SC, "vision_wenzhou.jpg"))
band = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(4.6),
                          prs.slide_width, Inches(2.9))
band.fill.solid()
band.fill.fore_color.rgb = BG
band.line.fill.background()
band.shadow.inherit = False
semi_transparent(band, 80000)
add_text(s, 0.8, 4.75, 12, 0.8, "让温州成为医疗数据要素流通的第一城", 34, WHITE,
         bold=True)
add_text(s, 0.8, 5.7, 12, 0.5,
         "430 项单元测试全绿 · 线上 Demo 每日可验 · 软著×3 申报中", 16, CYAN)
add_text(s, 0.8, 6.35, 12, 0.5, "瓯医数链 OuMedTrust · 恳请评委指正", 14, GREY)

prs.save(OUT_PPTX)
print("saved:", OUT_PPTX, "slides:", len(prs.slides._sldIdLst))
